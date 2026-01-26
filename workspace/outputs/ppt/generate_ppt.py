"""PPT 생성 스크립트 - 제안서 JSON 기반."""

import json
from datetime import datetime
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# 다크 테마 컬러
COLORS = {
    "background": "1E1E2E",
    "surface": "2D2D3F",
    "primary": "7C3AED",
    "secondary": "06B6D4",
    "accent": "F59E0B",
    "text_primary": "FFFFFF",
    "text_secondary": "A0AEC0",
    "success": "10B981",
    "warning": "F59E0B",
}


def hex_to_rgb(hex_color: str) -> RGBColor:
    hex_color = hex_color.lstrip('#')
    return RGBColor(
        int(hex_color[0:2], 16),
        int(hex_color[2:4], 16),
        int(hex_color[4:6], 16)
    )


def set_slide_background(slide, color_hex: str):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = hex_to_rgb(color_hex)


def add_title_slide(prs, title: str, subtitle: str = ""):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, COLORS["background"])

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(3), Inches(9), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = hex_to_rgb(COLORS["text_primary"])
    p.alignment = PP_ALIGN.CENTER

    if subtitle:
        subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.5), Inches(9), Inches(0.8))
        tf = subtitle_box.text_frame
        p = tf.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(20)
        p.font.color.rgb = hex_to_rgb(COLORS["text_secondary"])
        p.alignment = PP_ALIGN.CENTER

    return slide


def add_section_slide(prs, section_num: int, title: str):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, COLORS["background"])

    num_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1))
    tf = num_box.text_frame
    p = tf.paragraphs[0]
    p.text = f"0{section_num}" if section_num < 10 else str(section_num)
    p.font.size = Pt(72)
    p.font.bold = True
    p.font.color.rgb = hex_to_rgb(COLORS["primary"])
    p.alignment = PP_ALIGN.CENTER

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.5), Inches(9), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = hex_to_rgb(COLORS["text_primary"])
    p.alignment = PP_ALIGN.CENTER

    return slide


def add_content_slide(prs, title: str, bullets: list):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, COLORS["background"])

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = hex_to_rgb(COLORS["text_primary"])

    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(5))
    tf = content_box.text_frame
    tf.word_wrap = True

    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"• {bullet}"
        p.font.size = Pt(18)
        p.font.color.rgb = hex_to_rgb(COLORS["text_secondary"])
        p.space_before = Pt(10)

    return slide


def add_highlight_slide(prs, main_text: str, sub_text: str = ""):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, COLORS["background"])

    main_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(2))
    tf = main_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = main_text
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = hex_to_rgb(COLORS["text_primary"])
    p.alignment = PP_ALIGN.CENTER

    if sub_text:
        sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.5), Inches(9), Inches(1))
        tf = sub_box.text_frame
        p = tf.paragraphs[0]
        p.text = sub_text
        p.font.size = Pt(20)
        p.font.color.rgb = hex_to_rgb(COLORS["secondary"])
        p.alignment = PP_ALIGN.CENTER

    return slide


def add_two_column_slide(prs, title: str, left_title: str, left_items: list,
                          right_title: str, right_items: list):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, COLORS["background"])

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.7))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = hex_to_rgb(COLORS["text_primary"])

    # Left column
    left_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                       Inches(0.3), Inches(1.2), Inches(4.4), Inches(5.8))
    left_box.fill.solid()
    left_box.fill.fore_color.rgb = hex_to_rgb(COLORS["surface"])
    left_box.line.fill.background()

    left_title_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.4), Inches(4), Inches(0.5))
    tf = left_title_box.text_frame
    p = tf.paragraphs[0]
    p.text = left_title
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = hex_to_rgb(COLORS["warning"])

    left_content = slide.shapes.add_textbox(Inches(0.5), Inches(2.0), Inches(4), Inches(4.8))
    tf = left_content.text_frame
    tf.word_wrap = True
    for i, item in enumerate(left_items[:6]):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"• {item}"
        p.font.size = Pt(14)
        p.font.color.rgb = hex_to_rgb(COLORS["text_secondary"])
        p.space_before = Pt(6)

    # Right column
    right_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                        Inches(5.3), Inches(1.2), Inches(4.4), Inches(5.8))
    right_box.fill.solid()
    right_box.fill.fore_color.rgb = hex_to_rgb(COLORS["surface"])
    right_box.line.fill.background()

    right_title_box = slide.shapes.add_textbox(Inches(5.5), Inches(1.4), Inches(4), Inches(0.5))
    tf = right_title_box.text_frame
    p = tf.paragraphs[0]
    p.text = right_title
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = hex_to_rgb(COLORS["success"])

    right_content = slide.shapes.add_textbox(Inches(5.5), Inches(2.0), Inches(4), Inches(4.8))
    tf = right_content.text_frame
    tf.word_wrap = True
    for i, item in enumerate(right_items[:6]):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"• {item}"
        p.font.size = Pt(14)
        p.font.color.rgb = hex_to_rgb(COLORS["text_secondary"])
        p.space_before = Pt(6)

    return slide


def add_kpi_slide(prs, title: str, kpis: list):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, COLORS["background"])

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.7))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = hex_to_rgb(COLORS["text_primary"])

    positions = [(0.3, 1.3), (5.0, 1.3), (0.3, 4.0), (5.0, 4.0)]

    for i, kpi in enumerate(kpis[:4]):
        x, y = positions[i]

        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                       Inches(x), Inches(y), Inches(4.5), Inches(2.4))
        card.fill.solid()
        card.fill.fore_color.rgb = hex_to_rgb(COLORS["surface"])
        card.line.fill.background()

        metric_box = slide.shapes.add_textbox(Inches(x + 0.2), Inches(y + 0.2), Inches(4), Inches(0.4))
        tf = metric_box.text_frame
        p = tf.paragraphs[0]
        p.text = kpi.get("metric", "")
        p.font.size = Pt(16)
        p.font.color.rgb = hex_to_rgb(COLORS["text_secondary"])

        value_box = slide.shapes.add_textbox(Inches(x + 0.2), Inches(y + 0.7), Inches(4), Inches(0.8))
        tf = value_box.text_frame
        p = tf.paragraphs[0]
        p.text = kpi.get("value", "")
        p.font.size = Pt(24)
        p.font.bold = True
        p.font.color.rgb = hex_to_rgb(COLORS["primary"])

        improve_box = slide.shapes.add_textbox(Inches(x + 0.2), Inches(y + 1.6), Inches(4), Inches(0.5))
        tf = improve_box.text_frame
        p = tf.paragraphs[0]
        p.text = kpi.get("improvement", "")
        p.font.size = Pt(16)
        p.font.color.rgb = hex_to_rgb(COLORS["success"])

    return slide


def add_timeline_slide(prs, title: str, phases: list):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, COLORS["background"])

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.7))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = hex_to_rgb(COLORS["text_primary"])

    colors = [COLORS["primary"], COLORS["secondary"], COLORS["accent"], COLORS["success"], COLORS["warning"]]
    bar_y = 2.5

    for i, phase in enumerate(phases[:5]):
        y = bar_y + (i * 1.0)

        bar = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                      Inches(0.5), Inches(y), Inches(2.5), Inches(0.6))
        bar.fill.solid()
        bar.fill.fore_color.rgb = hex_to_rgb(colors[i % len(colors)])
        bar.line.fill.background()

        name_box = slide.shapes.add_textbox(Inches(0.6), Inches(y + 0.1), Inches(2.3), Inches(0.4))
        tf = name_box.text_frame
        p = tf.paragraphs[0]
        p.text = phase.get("phase_name", "")
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = hex_to_rgb(COLORS["text_primary"])

        duration_box = slide.shapes.add_textbox(Inches(3.2), Inches(y + 0.1), Inches(1.5), Inches(0.4))
        tf = duration_box.text_frame
        p = tf.paragraphs[0]
        p.text = phase.get("duration", "")
        p.font.size = Pt(14)
        p.font.color.rgb = hex_to_rgb(COLORS["secondary"])

        desc_box = slide.shapes.add_textbox(Inches(5.0), Inches(y + 0.1), Inches(4.5), Inches(0.4))
        tf = desc_box.text_frame
        p = tf.paragraphs[0]
        p.text = phase.get("description", "")
        p.font.size = Pt(12)
        p.font.color.rgb = hex_to_rgb(COLORS["text_secondary"])

    return slide


def add_team_slide(prs, title: str, team: list, total_mm: float):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, COLORS["background"])

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.7))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = hex_to_rgb(COLORS["text_primary"])

    for i, member in enumerate(team[:6]):
        row = i // 3
        col = i % 3
        x = 0.3 + (col * 3.2)
        y = 1.2 + (row * 2.6)

        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                       Inches(x), Inches(y), Inches(3.0), Inches(2.3))
        card.fill.solid()
        card.fill.fore_color.rgb = hex_to_rgb(COLORS["surface"])
        card.line.fill.background()

        role_box = slide.shapes.add_textbox(Inches(x + 0.15), Inches(y + 0.15), Inches(2.7), Inches(0.5))
        tf = role_box.text_frame
        p = tf.paragraphs[0]
        p.text = member.get("role", "")
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = hex_to_rgb(COLORS["primary"])

        count_box = slide.shapes.add_textbox(Inches(x + 0.15), Inches(y + 0.65), Inches(2.7), Inches(0.4))
        tf = count_box.text_frame
        p = tf.paragraphs[0]
        p.text = f"{member.get('count', 1)}명"
        p.font.size = Pt(22)
        p.font.bold = True
        p.font.color.rgb = hex_to_rgb(COLORS["text_primary"])

        resp = member.get("responsibilities", [])
        resp_text = ", ".join(resp[:2]) if resp else ""
        exp_box = slide.shapes.add_textbox(Inches(x + 0.15), Inches(y + 1.2), Inches(2.7), Inches(1.0))
        tf = exp_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = resp_text[:50]
        p.font.size = Pt(11)
        p.font.color.rgb = hex_to_rgb(COLORS["text_secondary"])

    total_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(9), Inches(0.5))
    tf = total_box.text_frame
    p = tf.paragraphs[0]
    p.text = f"총 투입 공수: {total_mm} Man-Months"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = hex_to_rgb(COLORS["secondary"])
    p.alignment = PP_ALIGN.CENTER

    return slide


def add_risk_slide(prs, title: str, risks: list):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, COLORS["background"])

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.7))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = hex_to_rgb(COLORS["text_primary"])

    y_start = 1.2
    for i, risk in enumerate(risks[:5]):
        y = y_start + (i * 1.1)

        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                      Inches(0.3), Inches(y), Inches(9.4), Inches(1.0))
        box.fill.solid()
        box.fill.fore_color.rgb = hex_to_rgb(COLORS["surface"])
        box.line.fill.background()

        level = risk.get("level", "MEDIUM")
        level_color = COLORS["warning"] if level == "HIGH" else COLORS["secondary"]
        level_box = slide.shapes.add_textbox(Inches(0.4), Inches(y + 0.1), Inches(1.0), Inches(0.4))
        tf = level_box.text_frame
        p = tf.paragraphs[0]
        p.text = level
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = hex_to_rgb(level_color)

        desc_box = slide.shapes.add_textbox(Inches(1.5), Inches(y + 0.1), Inches(4.0), Inches(0.4))
        tf = desc_box.text_frame
        p = tf.paragraphs[0]
        p.text = risk.get("description", "")[:45]
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = hex_to_rgb(COLORS["text_primary"])

        mit_box = slide.shapes.add_textbox(Inches(1.5), Inches(y + 0.5), Inches(8.0), Inches(0.4))
        tf = mit_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = f"→ {risk.get('mitigation', '')[:60]}"
        p.font.size = Pt(11)
        p.font.color.rgb = hex_to_rgb(COLORS["text_secondary"])

    return slide


def add_steps_slide(prs, title: str, steps: list):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, COLORS["background"])

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.7))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = hex_to_rgb(COLORS["text_primary"])

    y_start = 1.5
    for i, step in enumerate(steps[:5]):
        y = y_start + (i * 1.1)

        circle = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                         Inches(0.5), Inches(y), Inches(0.6), Inches(0.6))
        circle.fill.solid()
        circle.fill.fore_color.rgb = hex_to_rgb(COLORS["primary"])
        circle.line.fill.background()

        num_box = slide.shapes.add_textbox(Inches(0.5), Inches(y + 0.1), Inches(0.6), Inches(0.4))
        tf = num_box.text_frame
        p = tf.paragraphs[0]
        p.text = str(i + 1)
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = hex_to_rgb(COLORS["text_primary"])
        p.alignment = PP_ALIGN.CENTER

        action_box = slide.shapes.add_textbox(Inches(1.3), Inches(y + 0.1), Inches(8), Inches(0.5))
        tf = action_box.text_frame
        p = tf.paragraphs[0]
        p.text = step
        p.font.size = Pt(18)
        p.font.color.rgb = hex_to_rgb(COLORS["text_primary"])

    return slide


def add_closing_slide(prs, title: str = "감사합니다", subtitle: str = ""):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, COLORS["background"])

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(56)
    p.font.bold = True
    p.font.color.rgb = hex_to_rgb(COLORS["primary"])
    p.alignment = PP_ALIGN.CENTER

    if subtitle:
        subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.5), Inches(9), Inches(1))
        tf = subtitle_box.text_frame
        p = tf.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(20)
        p.font.color.rgb = hex_to_rgb(COLORS["text_secondary"])
        p.alignment = PP_ALIGN.CENTER

    return slide


def generate_ppt(data: dict, output_path: Path):
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    title = data.get("title", "프로젝트 제안서")
    client = data.get("client_name", "귀사")
    exec_summary = data.get("executive_summary", "")
    overview = data.get("project_overview", {})
    scope = data.get("scope_of_work", {})
    solution = data.get("solution_approach", {})
    timeline = data.get("timeline", {})
    resource = data.get("resource_plan", {})
    risks = data.get("risks", [])
    benefits = data.get("expected_benefits", [])
    next_steps = data.get("next_steps", [])
    metadata = data.get("metadata", {})

    # 1. 표지
    add_title_slide(prs, title.replace(" 제안서", ""), f"{metadata.get('created_at', '')[:10]} | {client} 귀중")

    # 2. 목차
    add_content_slide(prs, "목차 (Contents)", [
        "01  경영진 요약",
        "02  프로젝트 배경 및 목표",
        "03  작업 범위",
        "04  솔루션 및 기술 접근법",
        "05  일정 계획",
        "06  투입 인력",
        "07  리스크 관리",
        "08  기대 효과",
        "09  다음 단계"
    ])

    # 3. 경영진 요약 - 핵심 메시지
    add_highlight_slide(
        prs,
        "DOS에서 클라우드로, 수기에서 디지털로",
        "15년 된 레거시 시스템의 디지털 전환"
    )

    # 4. 경영진 요약 상세
    summary_lines = exec_summary.split('\n\n')[:3]
    add_content_slide(prs, "경영진 요약 (Executive Summary)", summary_lines)

    # 5. 현재 상황 섹션
    add_section_slide(prs, 1, "프로젝트 배경")

    # 6. 배경 및 문제점
    problems = [
        "15년 이상 운영된 DOS 기반 레거시 시스템",
        "유지보수 인력 확보의 어려움",
        "수기 대장 작성으로 인한 업무 비효율",
        "운전기사 평균 대기 시간 45분",
        "실시간 현황 파악 및 통계 분석 불가"
    ]
    add_content_slide(prs, "현재의 문제점", problems)

    # 7. Before vs After
    before = ["DOS 기반 레거시", "수기 전표 작성", "45분 대기 시간", "5% 데이터 오류율"]
    after = ["웹/모바일 시스템", "100% 전자 전표", "20분 이하 대기", "0.1% 이하 오류율"]
    add_two_column_slide(prs, "Before vs After", "현재 (AS-IS)", before, "미래 (TO-BE)", after)

    # 8. 프로젝트 목표
    objectives = overview.get("objectives", [])
    add_content_slide(prs, "프로젝트 목표", objectives[:6])

    # 9. 성공 기준 KPI
    success_criteria = overview.get("success_criteria", [])
    kpis = [
        {"metric": "대기 시간", "value": "45분 → 20분", "improvement": "55% 이상 단축"},
        {"metric": "전표 전자화", "value": "0% → 100%", "improvement": "완전 전자화"},
        {"metric": "데이터 오류율", "value": "5% → 0.1%", "improvement": "98% 감소"},
        {"metric": "시스템 가용성", "value": "99.5%+", "improvement": "클라우드 안정성"}
    ]
    add_kpi_slide(prs, "핵심 성공 지표 (KPI)", kpis)

    # 10. 작업 범위 섹션
    add_section_slide(prs, 2, "작업 범위")

    # 11. 포함/제외 범위
    in_scope = scope.get("in_scope", [])
    out_scope = scope.get("out_of_scope", [])
    add_two_column_slide(prs, "작업 범위 (Scope)", "포함 범위", in_scope[:6], "제외 범위", out_scope[:4])

    # 12. 주요 기능
    features = scope.get("key_features", [])
    feature_bullets = [f"{f.get('name', '')}: {f.get('description', '')}" for f in features]
    add_content_slide(prs, "주요 기능", feature_bullets)

    # 13. 솔루션 섹션
    add_section_slide(prs, 3, "솔루션 접근법")

    # 14. 솔루션 개요
    add_highlight_slide(
        prs,
        "클라우드 + 엣지 컴퓨팅 하이브리드",
        solution.get("overview", "")[:80]
    )

    # 15. 기술 스택
    tech_stack = solution.get("technology_stack", [])
    add_content_slide(prs, "기술 스택 (Technology Stack)", tech_stack[:6])

    # 16. 일정 섹션
    add_section_slide(prs, 4, "일정 계획")

    # 17. 타임라인
    phases = timeline.get("phases", [])
    total_duration = timeline.get("total_duration", "7개월")
    add_timeline_slide(prs, f"프로젝트 타임라인 ({total_duration})", phases)

    # 18. 팀 구성
    team = resource.get("team_structure", [])
    total_mm = resource.get("total_man_months", 14.5)
    add_team_slide(prs, "프로젝트 팀 구성", team, total_mm)

    # 19. 리스크 관리
    add_risk_slide(prs, "리스크 관리", risks)

    # 20. 기대 효과
    add_content_slide(prs, "기대 효과", benefits[:6])

    # 21. 다음 단계
    add_steps_slide(prs, "다음 단계 (Next Steps)", next_steps)

    # 22. 마무리
    add_closing_slide(prs, "감사합니다", "Q&A | 지금 바로 함께 시작하세요!")

    prs.save(str(output_path))
    return output_path


def main():
    print("\n" + "=" * 70)
    print("PPT 제안서 생성")
    print(f'시작 시간: {datetime.now().strftime("%Y-%m-%d %H:%M:%S")}')
    print("=" * 70)

    # 제안서 JSON 로드
    json_path = Path("workspace/outputs/proposals/PROP-20260126-100000.json")
    if not json_path.exists():
        print(f"제안서 파일을 찾을 수 없습니다: {json_path}")
        return

    with open(json_path, 'r', encoding='utf-8') as f:
        data = json.load(f)
    print(f"\n[입력] 제안서: {json_path}")

    # 출력 경로
    output_dir = Path("workspace/outputs/ppt")
    output_dir.mkdir(parents=True, exist_ok=True)

    timestamp = datetime.now().strftime("%Y%m%d-%H%M%S")
    output_path = output_dir / f"PPT-{timestamp}.pptx"

    # PPT 생성
    result = generate_ppt(data, output_path)

    if result:
        print(f"\n[완료] PPT 생성 완료: {output_path}")
        print(f"   슬라이드 수: 22장")
    else:
        print("\n[실패] PPT 생성 실패")


if __name__ == "__main__":
    main()
