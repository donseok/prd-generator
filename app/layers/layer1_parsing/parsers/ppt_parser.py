"""
파워포인트 파일(.pptx) 파서입니다.
python-pptx 라이브러리를 사용하여 슬라이드 텍스트, 표, 노트 내용을 추출합니다.
"""

from pathlib import Path
from typing import Optional

from app.models import InputType, ParsedContent, InputMetadata
from ..base_parser import BaseParser
from ..prompts.parsing_prompts import PPT_PARSING_PROMPT


class PPTParser(BaseParser):
    """프레젠테이션 파일을 처리하는 파서입니다."""

    @property
    def supported_types(self) -> list[InputType]:
        return [InputType.POWERPOINT]

    @property
    def supported_extensions(self) -> list[str]:
        return [".pptx", ".ppt"]

    async def parse(
        self,
        file_path: Path,
        metadata: Optional[dict] = None
    ) -> ParsedContent:
        """PPT 파일을 파싱하여 내용을 추출합니다."""
        from pptx import Presentation

        # 파일 로드
        prs = Presentation(str(file_path))

        # 슬라이드별 데이터 추출
        slides_data = []
        for i, slide in enumerate(prs.slides, 1):
            slide_info = self._extract_slide_content(slide, i)
            slides_data.append(slide_info)

        # 텍스트 형태의 통합 본문 생성
        raw_text = self._build_raw_text(slides_data)

        # 메타데이터 생성
        file_metadata = await self.extract_metadata(file_path)
        file_metadata.slide_count = len(prs.slides)

        # 슬라이드별 섹션 생성
        sections = [
            {
                "title": f"슬라이드 {s['number']}: {s['title']}",
                "content": s["content"],
            }
            for s in slides_data
        ]

        # 구조 데이터 생성
        structured_data = {
            "slide_count": len(prs.slides),
            "slides": slides_data,
        }

        # AI(Claude) 분석 (가능한 경우)
        if self.claude_client:
            try:
                analysis = await self._analyze_with_claude(raw_text)
                if analysis and isinstance(analysis, dict) and analysis:
                    structured_data["ai_analysis"] = analysis
                else:
                    # AI 분석 실패 시 슬라이드 내용에서 직접 분석 데이터 생성
                    structured_data["ai_analysis"] = self._create_fallback_analysis(slides_data)
            except Exception as e:
                print(f"Claude PPT 분석 실패: {e}")
                structured_data["ai_analysis"] = self._create_fallback_analysis(slides_data)

        return ParsedContent(
            raw_text=raw_text,
            structured_data=structured_data,
            metadata=file_metadata,
            sections=sections,
        )

    def _extract_slide_content(self, slide, slide_number: int) -> dict:
        """단일 슬라이드에서 텍스트와 표 내용을 추출합니다."""
        title = ""
        content_parts = []
        notes = ""

        for shape in slide.shapes:
            # 텍스트 상자 처리
            if shape.has_text_frame:
                text = shape.text_frame.text.strip()
                if text:
                    # 제목 추정: 플레이스홀더 타입이 제목인 경우
                    if not title and shape.is_placeholder:
                        if hasattr(shape, 'placeholder_format'):
                            ph_type = shape.placeholder_format.type
                            # TITLE(1), CENTER_TITLE(3) 타입
                            if ph_type in [1, 3]:
                                title = text
                                continue
                    content_parts.append(text)

            # 표 처리
            if shape.has_table:
                table_text = self._extract_table(shape.table)
                content_parts.append(table_text)

        # 발표자 노트 추출
        if slide.has_notes_slide:
            notes_frame = slide.notes_slide.notes_text_frame
            if notes_frame:
                notes = notes_frame.text.strip()

        return {
            "number": slide_number,
            "title": title or f"슬라이드 {slide_number}",
            "content": "\n".join(content_parts),
            "notes": notes,
            "has_images": any(
                shape.shape_type == 13  # 그림 타입
                for shape in slide.shapes
                if hasattr(shape, 'shape_type')
            ),
        }

    def _extract_table(self, table) -> str:
        """표 내용을 텍스트(파이프 | 구분)로 변환합니다."""
        rows = []
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells]
            rows.append(" | ".join(cells))
        return "\n".join(rows)

    def _build_raw_text(self, slides_data: list) -> str:
        """모든 슬라이드 내용을 하나의 텍스트로 합칩니다."""
        parts = []

        for slide in slides_data:
            parts.append(f"=== 슬라이드 {slide['number']}: {slide['title']} ===")
            if slide['content']:
                parts.append(slide['content'])
            if slide['notes']:
                parts.append(f"\n[발표자 노트]\n{slide['notes']}")
            parts.append("")

        return "\n".join(parts)

    async def _analyze_with_claude(self, raw_text: str) -> dict:
        """Claude에게 PPT 내용 분석을 요청합니다."""
        result = await self.claude_client.complete_json(
            system_prompt=PPT_PARSING_PROMPT,
            user_prompt=f"다음 PPT 내용을 분석해주세요:\n\n{raw_text[:8000]}",
            temperature=0.2,
        )
        return result

    def _create_fallback_analysis(self, slides_data: list) -> dict:
        """AI 분석 실패 시 슬라이드 데이터에서 기본 분석 결과를 생성합니다."""
        # 슬라이드 제목들을 수집
        titles = [s.get('title', '') for s in slides_data if s.get('title')]

        # 모든 내용을 합침
        all_content = "\n".join([s.get('content', '') for s in slides_data])

        # 키워드 기반 주제 감지
        topics = []
        keyword_topics = {
            "시스템": "시스템 구축",
            "플랫폼": "플랫폼 개발",
            "데이터": "데이터 관리",
            "사용자": "사용자 경험",
            "보안": "보안 요구사항",
            "성능": "성능 최적화",
            "통합": "시스템 통합",
            "API": "API 개발",
            "모바일": "모바일 앱",
            "웹": "웹 서비스",
            "관리": "관리 시스템",
            "자동화": "프로세스 자동화",
            "계량": "계량 시스템",
            "IoT": "IoT 연동",
            "스마트": "스마트 시스템",
        }

        for keyword, topic in keyword_topics.items():
            if keyword in all_content or any(keyword in t for t in titles):
                topics.append(topic)

        if not topics:
            topics = ["시스템 구축 프로젝트"]

        # 기본 요구사항 추출 (슬라이드 제목 기반)
        requirements = []
        for slide in slides_data:
            title = slide.get('title', '')
            content = slide.get('content', '')

            # 제목이 있고 내용이 있는 슬라이드를 요구사항으로 변환
            if title and content and len(content) > 20:
                # "슬라이드 X:" 형식 제거
                clean_title = title
                if "슬라이드" in clean_title:
                    parts = clean_title.split(":", 1)
                    if len(parts) > 1:
                        clean_title = parts[1].strip()

                if clean_title and clean_title != f"슬라이드 {slide.get('number', '')}":
                    requirements.append({
                        "title": clean_title[:100],
                        "description": content[:500],
                        "source": f"슬라이드 {slide.get('number', '')}",
                    })

        return {
            "document_type": "프로젝트 수행계획서",
            "main_topics": topics[:5],
            "key_requirements": requirements[:15],
            "stakeholders": ["프로젝트 관리자", "개발팀", "고객사"],
            "estimated_complexity": "중간" if len(slides_data) < 30 else "높음",
            "analysis_source": "슬라이드 내용 직접 분석",
        }