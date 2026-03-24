"""PPT 제안서 생성 스크립트 (2026 Premium Design).

Usage:
    python -m app.scripts.ppt_maker
    python -m app.scripts.ppt_maker --theme modern-blue
    python -m app.scripts.ppt_maker --theme dark-corporate

제안서(PROP-*.json)를 기반으로 프리미엄 PPT 생성.
16:9 와이드스크린, 다양한 레이아웃, 데이터 시각화, 그라데이션 배경.

슬라이드 안전 영역: 높이 7.5인치, 콘텐츠는 0.3~7.2인치 이내.
"""

import argparse
import json
import re
import math
from datetime import datetime
from io import BytesIO
from pathlib import Path

import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt
import matplotlib.patches as mpatches
import matplotlib.font_manager as fm
from PIL import Image, ImageDraw

# 한글 폰트 설정 (macOS: Apple SD Gothic Neo, Windows: 맑은 고딕, fallback: Nanum Gothic)
_KR_FONT_CANDIDATES = ['Apple SD Gothic Neo', 'Malgun Gothic', 'Nanum Gothic', 'NanumGothic']
_KR_FONT = None
for _fname in _KR_FONT_CANDIDATES:
    if any(f.name == _fname for f in fm.fontManager.ttflist):
        _KR_FONT = _fname
        break
if _KR_FONT:
    plt.rcParams['font.family'] = _KR_FONT
plt.rcParams['axes.unicode_minus'] = False

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

from app.scripts.arch_diagram import generate_from_trd_file


# ============================================================
# Theme System
# ============================================================
PPT_THEMES = {
    "modern-blue": {
        "background": "FFFFFF",
        "bg_light": "F8FAFC",
        "bg_section": "1E3A5F",
        "card": "F1F5F9",
        "card_alt": "EEF2FF",
        "primary": "2563EB",
        "secondary": "7C3AED",
        "accent3": "0891B2",
        "accent4": "059669",
        "accent5": "DC2626",
        "orange": "EA580C",
        "amber": "D97706",
        "title": "0F172A",
        "body": "334155",
        "subtle": "64748B",
        "muted": "94A3B8",
        "border": "E2E8F0",
        "white": "FFFFFF",
        "gradient_start": "1E3A5F",
        "gradient_end": "2D4A6F",
        "shadow": "CBD5E1",
        "card_before": "FEF2F2",
        "card_after": "F0FDF4",
        "font_title": "맑은 고딕",
        "font_body": "맑은 고딕",
    },
    "dark-corporate": {
        "background": "1A1A2E",
        "bg_light": "16213E",
        "bg_section": "0F3460",
        "card": "16213E",
        "card_alt": "1A1A3E",
        "primary": "E94560",
        "secondary": "533483",
        "accent3": "0891B2",
        "accent4": "10B981",
        "accent5": "EF4444",
        "orange": "F97316",
        "amber": "F59E0B",
        "title": "F1F5F9",
        "body": "CBD5E1",
        "subtle": "94A3B8",
        "muted": "64748B",
        "border": "334155",
        "white": "F8FAFC",
        "gradient_start": "0F3460",
        "gradient_end": "1A1A4E",
        "shadow": "0F172A",
        "card_before": "3B1A1A",
        "card_after": "1A3B1A",
        "font_title": "맑은 고딕",
        "font_body": "맑은 고딕",
    },
}

# Current theme (set by CLI or default)
_current_theme = PPT_THEMES["modern-blue"]


def _t(key: str) -> str:
    """Get color from current theme."""
    return _current_theme.get(key, "334155")


def _palette() -> list:
    """Get palette colors from current theme."""
    return [_t("primary"), _t("secondary"), _t("accent3"), _t("accent4"), _t("orange"), _t("amber")]


# ============================================================
# 16:9 Widescreen Layout Constants
# ============================================================
SLIDE_W = 13.333
SLIDE_H = 7.5
MARGIN_X = 0.7
MARGIN_TOP = 0.4
CONTENT_TOP = 1.8
SAFE_BOTTOM = 7.1
CONTENT_H = SAFE_BOTTOM - CONTENT_TOP
CONTENT_W = SLIDE_W - 2 * MARGIN_X

# Font size constants
FONT_SIZE_TITLE = 32
FONT_SIZE_SUBTITLE = 18
FONT_SIZE_BODY = 13
FONT_SIZE_SMALL = 11
FONT_SIZE_METRIC = 48
FONT_SIZE_METRIC_LABEL = 12


# ============================================================
# Utility Functions
# ============================================================

def hex_to_rgb(hex_color: str) -> RGBColor:
    h = hex_color.lstrip('#')
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def hex_to_tuple(hex_color: str) -> tuple:
    """Convert hex color to (r, g, b) tuple for matplotlib."""
    h = hex_color.lstrip('#')
    return (int(h[0:2], 16) / 255, int(h[2:4], 16) / 255, int(h[4:6], 16) / 255)


def set_slide_bg(slide, color_hex: str):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = hex_to_rgb(color_hex)


def _add_text(slide, left, top, width, height, text, font_size=18,
              bold=False, color=None, alignment=PP_ALIGN.LEFT,
              font_name=None, word_wrap=True):
    if color is None:
        color = _t("body")
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = word_wrap
    p = tf.paragraphs[0]
    p.text = str(text)
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = hex_to_rgb(color)
    p.font.name = font_name or _current_theme.get("font_body", "맑은 고딕")
    p.alignment = alignment
    return box


def _add_rounded_card(slide, left, top, width, height, fill_color=None):
    if fill_color is None:
        fill_color = _t("card")
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = hex_to_rgb(fill_color)
    shape.line.fill.background()
    shape.adjustments[0] = 0.06
    return shape


def _add_card_with_shadow(slide, left, top, width, height, fill_color=None, shadow_offset=0.05):
    """Add a card with a simulated drop shadow behind it."""
    if fill_color is None:
        fill_color = _t("card")
    # Shadow (slightly offset, lighter color)
    shadow = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left + shadow_offset), Inches(top + shadow_offset),
        Inches(width), Inches(height)
    )
    shadow.fill.solid()
    shadow.fill.fore_color.rgb = hex_to_rgb(_t("shadow"))
    shadow.line.fill.background()
    shadow.adjustments[0] = 0.06

    # Main card on top
    card = _add_rounded_card(slide, left, top, width, height, fill_color)
    return card


def _add_circle_icon(slide, left, top, size, color, text="", text_color=None):
    if text_color is None:
        text_color = _t("white")
    circle = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(left), Inches(top), Inches(size), Inches(size)
    )
    circle.fill.solid()
    circle.fill.fore_color.rgb = hex_to_rgb(color)
    circle.line.fill.background()
    if text:
        tf = circle.text_frame
        tf.word_wrap = False
        p = tf.paragraphs[0]
        p.text = str(text)
        p.font.size = Pt(max(9, int(size * 13)))
        p.font.bold = True
        p.font.color.rgb = hex_to_rgb(text_color)
        p.alignment = PP_ALIGN.CENTER
    return circle


def _add_accent_bar(slide, left, top, width, height=0.05, color=None):
    if color is None:
        color = _t("primary")
    bar = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = hex_to_rgb(color)
    bar.line.fill.background()
    return bar


def _add_vertical_line(slide, left, top, height, width=0.03, color=None):
    """Add a thin vertical divider line."""
    if color is None:
        color = _t("border")
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = hex_to_rgb(color)
    line.line.fill.background()
    return line


def _add_progress_bar(slide, left, top, width, height, pct, bg_color=None, fill_color=None):
    if bg_color is None:
        bg_color = _t("border")
    if fill_color is None:
        fill_color = _t("primary")
    bg = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = hex_to_rgb(bg_color)
    bg.line.fill.background()
    fill_w = max(width * (pct / 100.0), 0.1)
    fg = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left), Inches(top), Inches(fill_w), Inches(height)
    )
    fg.fill.solid()
    fg.fill.fore_color.rgb = hex_to_rgb(fill_color)
    fg.line.fill.background()
    return fg


def _slide_header(slide, eng_label, kor_title, label_color=None, bar_color=None):
    """공통 슬라이드 헤더 (와이드스크린 대응)."""
    if label_color is None:
        label_color = _t("primary")
    if bar_color is None:
        bar_color = _t("primary")
    _add_text(slide, 0.8, 0.4, 10, 0.4, eng_label,
              font_size=FONT_SIZE_METRIC_LABEL, bold=True, color=label_color)
    _add_text(slide, 0.8, 0.8, 10, 0.6, kor_title,
              font_size=FONT_SIZE_TITLE, bold=True, color=_t("title"),
              font_name=_current_theme.get("font_title", "맑은 고딕"))
    _add_accent_bar(slide, 0.8, 1.4, 1.5, 0.04, bar_color)


def _calc_item_spacing(n_items, start_y=CONTENT_TOP, end_y=SAFE_BOTTOM, item_h=0.8):
    """항목 수에 따라 동적 간격 계산."""
    available = end_y - start_y
    if n_items <= 0:
        return start_y, item_h
    total_needed = n_items * item_h
    if total_needed > available:
        item_h = available / n_items
    return start_y, item_h


# ============================================================
# Gradient Background Generator (Pillow)
# ============================================================

def _create_gradient_image(width, height, color1, color2, direction='horizontal'):
    """Create gradient background image using Pillow. Returns BytesIO PNG."""
    img = Image.new('RGB', (width, height))
    draw = ImageDraw.Draw(img)

    r1, g1, b1 = int(color1[0:2], 16), int(color1[2:4], 16), int(color1[4:6], 16)
    r2, g2, b2 = int(color2[0:2], 16), int(color2[2:4], 16), int(color2[4:6], 16)

    if direction == 'horizontal':
        for x in range(width):
            ratio = x / max(width - 1, 1)
            r = int(r1 + (r2 - r1) * ratio)
            g = int(g1 + (g2 - g1) * ratio)
            b = int(b1 + (b2 - b1) * ratio)
            draw.line([(x, 0), (x, height)], fill=(r, g, b))
    elif direction == 'vertical':
        for y in range(height):
            ratio = y / max(height - 1, 1)
            r = int(r1 + (r2 - r1) * ratio)
            g = int(g1 + (g2 - g1) * ratio)
            b = int(b1 + (b2 - b1) * ratio)
            draw.line([(0, y), (width, y)], fill=(r, g, b))
    elif direction == 'diagonal':
        for y in range(height):
            for x in range(width):
                ratio = (x / max(width - 1, 1) + y / max(height - 1, 1)) / 2
                r = int(r1 + (r2 - r1) * ratio)
                g = int(g1 + (g2 - g1) * ratio)
                b = int(b1 + (b2 - b1) * ratio)
                draw.point((x, y), fill=(r, g, b))

    buf = BytesIO()
    img.save(buf, format='PNG')
    buf.seek(0)
    return buf


def _set_gradient_bg(slide, color1, color2, direction='diagonal'):
    """Set a gradient background on a slide using a full-slide image."""
    w_px = 1920
    h_px = 1080
    buf = _create_gradient_image(w_px, h_px, color1, color2, direction)
    slide.shapes.add_picture(buf, Inches(0), Inches(0), Inches(SLIDE_W), Inches(SLIDE_H))


# ============================================================
# Chart Generation (matplotlib)
# ============================================================

def _create_donut_chart(data, colors, title=""):
    """Create a donut chart and return as BytesIO PNG.
    data: list of (label, value) tuples
    colors: list of hex color strings
    """
    if not data:
        return None

    labels = [d[0] for d in data]
    values = [max(float(d[1]), 0.01) for d in data]
    chart_colors = [hex_to_tuple(c) for c in colors[:len(values)]]

    fig, ax = plt.subplots(figsize=(4, 4), facecolor='none')
    wedges, texts, autotexts = ax.pie(
        values, labels=None, colors=chart_colors,
        autopct='%1.0f%%', startangle=90, pctdistance=0.78,
        wedgeprops=dict(width=0.35, edgecolor='white', linewidth=2)
    )
    for t in autotexts:
        t.set_fontsize(10)
        t.set_fontweight('bold')
        t.set_color('white')

    # Legend
    ax.legend(wedges, labels, loc='center', fontsize=8, frameon=False,
              bbox_to_anchor=(0.5, 0.5))
    if title:
        ax.set_title(title, fontsize=12, fontweight='bold', pad=10)
    ax.set_aspect('equal')

    buf = BytesIO()
    fig.savefig(buf, format='png', dpi=150, transparent=True, bbox_inches='tight')
    plt.close(fig)
    buf.seek(0)
    return buf


def _create_bar_chart(categories, values, colors, title="", horizontal=True):
    """Create a bar chart and return as BytesIO PNG.
    categories: list of label strings
    values: list of numeric values
    colors: list of hex color strings
    """
    if not categories or not values:
        return None

    n = len(categories)
    chart_colors = [hex_to_tuple(colors[i % len(colors)]) for i in range(n)]

    fig, ax = plt.subplots(figsize=(6, max(2.5, n * 0.6)), facecolor='none')
    ax.set_facecolor('none')

    if horizontal:
        bars = ax.barh(range(n), values, color=chart_colors, height=0.6, edgecolor='none')
        ax.set_yticks(range(n))
        ax.set_yticklabels(categories, fontsize=9)
        ax.invert_yaxis()
        for bar, val in zip(bars, values):
            ax.text(bar.get_width() + max(values) * 0.02, bar.get_y() + bar.get_height() / 2,
                    f'{val}', va='center', fontsize=9, fontweight='bold')
    else:
        bars = ax.bar(range(n), values, color=chart_colors, width=0.6, edgecolor='none')
        ax.set_xticks(range(n))
        ax.set_xticklabels(categories, fontsize=9, rotation=15, ha='right')
        for bar, val in zip(bars, values):
            ax.text(bar.get_x() + bar.get_width() / 2, bar.get_height() + max(values) * 0.02,
                    f'{val}', ha='center', fontsize=9, fontweight='bold')

    if title:
        ax.set_title(title, fontsize=11, fontweight='bold', pad=8)

    ax.spines['top'].set_visible(False)
    ax.spines['right'].set_visible(False)
    ax.spines['left'].set_visible(False)
    ax.spines['bottom'].set_color('#E2E8F0')
    ax.tick_params(left=False, bottom=False)
    ax.grid(axis='x' if horizontal else 'y', alpha=0.15)

    buf = BytesIO()
    fig.savefig(buf, format='png', dpi=150, transparent=True, bbox_inches='tight')
    plt.close(fig)
    buf.seek(0)
    return buf


def _create_timeline_chart(phases, durations):
    """Create a Gantt-chart style timeline and return as BytesIO PNG.
    phases: list of phase name strings
    durations: list of duration strings or numeric values
    """
    if not phases:
        return None

    n = len(phases)
    # Parse durations to relative widths
    widths = []
    for d in durations:
        # Try to extract number from duration string
        nums = re.findall(r'(\d+)', str(d))
        widths.append(int(nums[0]) if nums else 1)

    palette = _palette()
    chart_colors = [hex_to_tuple(palette[i % len(palette)]) for i in range(n)]

    fig, ax = plt.subplots(figsize=(8, max(2, n * 0.55)), facecolor='none')
    ax.set_facecolor('none')

    cumulative = 0
    for i, (phase, w) in enumerate(zip(phases, widths)):
        ax.barh(i, w, left=cumulative, height=0.5, color=chart_colors[i],
                edgecolor='white', linewidth=1, alpha=0.9)
        # Phase label inside bar
        ax.text(cumulative + w / 2, i, phase[:20], ha='center', va='center',
                fontsize=8, fontweight='bold', color='white')
        # Duration label
        ax.text(cumulative + w + 0.1, i, str(durations[i])[:15],
                ha='left', va='center', fontsize=7, color='#64748B')
        cumulative += w

    ax.set_yticks([])
    ax.set_xticks([])
    ax.invert_yaxis()
    for spine in ax.spines.values():
        spine.set_visible(False)

    buf = BytesIO()
    fig.savefig(buf, format='png', dpi=150, transparent=True, bbox_inches='tight')
    plt.close(fig)
    buf.seek(0)
    return buf


def _create_gauge_chart(value, max_value, label=""):
    """Create a semi-circle gauge chart and return as BytesIO PNG."""
    if max_value <= 0:
        return None

    pct = min(value / max_value, 1.0)
    fig, ax = plt.subplots(figsize=(3, 2), facecolor='none', subplot_kw={'aspect': 'equal'})

    # Background arc
    theta1, theta2 = 0, 180
    bg_wedge = mpatches.Wedge((0.5, 0), 0.4, theta1, theta2, width=0.12,
                               facecolor='#E2E8F0', edgecolor='none',
                               transform=ax.transAxes)
    ax.add_patch(bg_wedge)

    # Value arc - color based on percentage
    if pct < 0.33:
        color = hex_to_tuple(_t("accent5"))
    elif pct < 0.66:
        color = hex_to_tuple(_t("amber"))
    else:
        color = hex_to_tuple(_t("accent4"))

    value_angle = theta1 + (theta2 - theta1) * pct
    val_wedge = mpatches.Wedge((0.5, 0), 0.4, theta1, value_angle, width=0.12,
                                facecolor=color, edgecolor='none',
                                transform=ax.transAxes)
    ax.add_patch(val_wedge)

    # Center text
    ax.text(0.5, 0.15, f'{int(pct * 100)}%', ha='center', va='center',
            fontsize=18, fontweight='bold', color=color, transform=ax.transAxes)
    if label:
        ax.text(0.5, -0.05, label, ha='center', va='center',
                fontsize=8, color='#64748B', transform=ax.transAxes)

    ax.set_xlim(0, 1)
    ax.set_ylim(-0.1, 0.6)
    ax.axis('off')

    buf = BytesIO()
    fig.savefig(buf, format='png', dpi=150, transparent=True, bbox_inches='tight')
    plt.close(fig)
    buf.seek(0)
    return buf


def _insert_chart(slide, chart_buf, left, top, width, height):
    """Insert a chart image (BytesIO) into a slide, fitting within bounds."""
    if chart_buf is None:
        return
    slide.shapes.add_picture(chart_buf, Inches(left), Inches(top), Inches(width), Inches(height))


# ============================================================
# Slide Builders
# ============================================================

def slide_cover(prs, title, subtitle, date_str):
    """슬라이드 1: 표지 (그라데이션 배경)."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    # Gradient background for cover
    _set_gradient_bg(slide, _t("gradient_start"), _t("gradient_end"), 'diagonal')

    # Top accent bar
    _add_accent_bar(slide, 0, 0, SLIDE_W, 0.07, _t("primary"))

    # Left vertical accent
    block = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0.07), Inches(0.22), Inches(7.36)
    )
    block.fill.solid()
    block.fill.fore_color.rgb = hex_to_rgb(_t("primary"))
    block.line.fill.background()

    # Title - large, centered
    _add_text(slide, 1.0, 2.0, 11, 1.2, title,
              font_size=46, bold=True, color=_t("white"),
              font_name=_current_theme.get("font_title", "맑은 고딕"))
    _add_accent_bar(slide, 1.0, 3.4, 2.5, 0.05, _t("primary"))

    if subtitle:
        _add_text(slide, 1.0, 3.7, 11, 0.5, subtitle,
                  font_size=20, color=_t("muted"))
    _add_text(slide, 1.0, 4.4, 11, 0.4, date_str,
              font_size=15, color=_t("muted"))

    # Bottom accent bar
    _add_accent_bar(slide, 0, 7.43, SLIDE_W, 0.07, _t("secondary"))


def slide_toc(prs, sections):
    """슬라이드 2: 목차 (3컬럼 와이드스크린 대응)."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, _t("background"))
    _slide_header(slide, "CONTENTS", "목차")

    n = len(sections)
    cols = 3 if n > 8 else 2
    col_count = math.ceil(n / cols)
    col_w = CONTENT_W / cols
    palette = _palette()

    for i, section in enumerate(sections):
        col = i // col_count
        row = i % col_count
        x_base = MARGIN_X + col * col_w
        spacing = min(0.45, (SAFE_BOTTOM - CONTENT_TOP) / max(col_count, 1))
        y = CONTENT_TOP + row * spacing

        num = f"{i + 1:02d}"
        color = palette[i % len(palette)]
        _add_circle_icon(slide, x_base, y, 0.36, color, num)
        _add_text(slide, x_base + 0.5, y + 0.01, col_w - 0.7, 0.36, section,
                  font_size=14, color=_t("body"))


def slide_exec_highlight(prs, main_text, metrics):
    """슬라이드 3: 경영진 요약 핵심 (Full-width highlight layout)."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    # Gradient section background
    _set_gradient_bg(slide, _t("bg_section"), _t("gradient_end"), 'diagonal')

    _add_text(slide, 1.0, 0.6, 11, 0.4, "EXECUTIVE SUMMARY",
              font_size=FONT_SIZE_METRIC_LABEL, bold=True, color=_t("muted"))

    # Main message - large centered text
    _add_text(slide, 1.0, 1.2, 11.3, 2.0, main_text[:180],
              font_size=26, bold=True, color=_t("white"),
              font_name=_current_theme.get("font_title", "맑은 고딕"))

    if metrics:
        n = min(len(metrics), 3)
        card_w = 3.2
        gap = 0.4
        total_w = n * card_w + (n - 1) * gap
        start_x = (SLIDE_W - total_w) / 2

        for i, m in enumerate(metrics[:3]):
            x = start_x + i * (card_w + gap)
            _add_card_with_shadow(slide, x, 3.8, card_w, 2.2,
                                  _t("gradient_end"), 0.06)

            # Label - small above
            _add_text(slide, x + 0.2, 3.95, card_w - 0.4, 0.35,
                      m.get("label", ""), font_size=FONT_SIZE_SMALL,
                      color=_t("muted"))
            # Value - oversized metric font
            _add_text(slide, x + 0.2, 4.4, card_w - 0.4, 0.8,
                      m.get("value", ""), font_size=FONT_SIZE_METRIC,
                      bold=True, color=_t("white"),
                      alignment=PP_ALIGN.CENTER,
                      font_name=_current_theme.get("font_title", "맑은 고딕"))
            if m.get("desc"):
                _add_text(slide, x + 0.2, 5.4, card_w - 0.4, 0.35,
                          m["desc"], font_size=10, color=_t("muted"),
                          alignment=PP_ALIGN.CENTER)


def slide_exec_detail(prs, items):
    """슬라이드 4: 경영진 요약 상세 (Quote/callout layout)."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, _t("background"))
    _slide_header(slide, "EXECUTIVE SUMMARY", "프로젝트 개요")

    if not items:
        return

    # Large quotation mark graphic
    _add_text(slide, 0.6, 1.5, 1.5, 1.2, "\u201C",
              font_size=72, bold=True, color=_t("border"),
              font_name=_current_theme.get("font_title", "맑은 고딕"))

    n = min(len(items), 5)
    start_y = CONTENT_TOP + 0.1
    available_h = SAFE_BOTTOM - start_y
    item_h = min(0.9, available_h / max(n, 1))

    for i, item in enumerate(items[:n]):
        y = start_y + i * item_h

        # Accent bar on left
        _add_accent_bar(slide, 1.2, y + 0.05, 0.06, item_h - 0.15,
                        _palette()[i % len(_palette())])

        # Label - bold accent color
        _add_text(slide, 1.5, y + 0.02, 2.8, 0.3, item.get("label", ""),
                  font_size=13, bold=True, color=_t("primary"))
        # Value - body text
        _add_text(slide, 1.5, y + 0.35, 10.5, 0.4, item.get("value", "")[:100],
                  font_size=15, color=_t("body"))

    # Closing quotation mark
    _add_text(slide, 11.5, SAFE_BOTTOM - 0.8, 1.5, 1.0, "\u201D",
              font_size=72, bold=True, color=_t("border"),
              font_name=_current_theme.get("font_title", "맑은 고딕"))


def slide_section(prs, num, title, subtitle=""):
    """섹션 구분 슬라이드 (그라데이션 배경)."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_gradient_bg(slide, _t("bg_section"), _t("gradient_end"), 'horizontal')

    _add_text(slide, 1.0, 2.0, 11, 0.9,
              f"{num:02d}" if isinstance(num, int) else str(num),
              font_size=72, bold=True, color=_t("primary"),
              font_name=_current_theme.get("font_title", "맑은 고딕"))
    _add_accent_bar(slide, 1.0, 3.1, 2.5, 0.05, _t("primary"))
    _add_text(slide, 1.0, 3.4, 11, 0.9, title,
              font_size=40, bold=True, color=_t("white"),
              font_name=_current_theme.get("font_title", "맑은 고딕"))
    if subtitle:
        _add_text(slide, 1.0, 4.3, 11, 0.5, subtitle,
                  font_size=17, color=_t("muted"))


def slide_challenges(prs, title, challenges):
    """도전과제 슬라이드 (2x2 카드 그리드, 와이드스크린)."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, _t("background"))
    _slide_header(slide, "CHALLENGES", title, _t("accent5"), _t("accent5"))

    n = min(len(challenges), 4)
    cols = 2
    rows = math.ceil(n / cols)
    cw = 5.7
    gap_x = 0.5
    gap_y = 0.3
    ch = (SAFE_BOTTOM - CONTENT_TOP - (rows - 1) * gap_y) / rows
    ch = min(ch, 2.4)
    start_x = (SLIDE_W - (cols * cw + (cols - 1) * gap_x)) / 2

    warn_icons = ["!", "!!", "?", "X"]
    for i, c in enumerate(challenges[:n]):
        col_idx = i % cols
        row_idx = i // cols
        x = start_x + col_idx * (cw + gap_x)
        y = CONTENT_TOP + row_idx * (ch + gap_y)

        _add_card_with_shadow(slide, x, y, cw, ch, _t("card"))
        _add_circle_icon(slide, x + 0.2, y + 0.2, 0.4, _t("accent5"),
                         warn_icons[i % len(warn_icons)])

        area = c.get("area", c.get("title", f"과제 {i + 1}"))
        desc = c.get("symptom", c.get("description", c.get("issue", "")))

        _add_text(slide, x + 0.75, y + 0.2, cw - 1.0, 0.35, area[:35],
                  font_size=15, bold=True, color=_t("title"))
        _add_text(slide, x + 0.2, y + 0.7, cw - 0.4, ch - 0.85, desc[:120],
                  font_size=FONT_SIZE_BODY, color=_t("body"))


def slide_risks_no_change(prs, risks):
    """변화하지 않으면 슬라이드."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, _t("background"))
    _slide_header(slide, "IF NO CHANGE", "변화하지 않으면?", _t("accent5"), _t("accent5"))

    n = min(len(risks), 5)
    start_y, item_h = _calc_item_spacing(n, CONTENT_TOP, SAFE_BOTTOM, 0.85)

    for i, risk in enumerate(risks[:n]):
        y = start_y + i * item_h
        text = risk if isinstance(risk, str) else str(risk)
        _add_card_with_shadow(slide, MARGIN_X, y, CONTENT_W, item_h - 0.08, _t("card"))
        _add_circle_icon(slide, MARGIN_X + 0.15, y + (item_h - 0.48) / 2 - 0.04, 0.42,
                         _t("accent5"), str(i + 1))
        _add_text(slide, MARGIN_X + 0.75, y + (item_h - 0.4) / 2 - 0.04, CONTENT_W - 1.0, 0.4,
                  text[:120], font_size=15, color=_t("body"))


def slide_before_after(prs, before_items, after_items):
    """Before vs After 비교 (Split 60/40 layout)."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, _t("background"))
    _slide_header(slide, "COMPARISON", "Before vs After")

    card_top = CONTENT_TOP
    card_h = SAFE_BOTTOM - card_top

    # Split: 55% left, 45% right with divider
    left_w = CONTENT_W * 0.48
    right_w = CONTENT_W * 0.48
    left_x = MARGIN_X
    divider_x = MARGIN_X + left_w + CONTENT_W * 0.02
    right_x = divider_x + CONTENT_W * 0.02

    # Left card: Before (red tinted)
    _add_card_with_shadow(slide, left_x, card_top, left_w, card_h, _t("card_before"))
    _add_circle_icon(slide, left_x + 0.15, card_top + 0.15, 0.45, _t("accent5"), "X")
    _add_text(slide, left_x + 0.7, card_top + 0.18, left_w - 1.0, 0.4, "AS-IS (현재)",
              font_size=18, bold=True, color=_t("accent5"))

    n_before = min(len(before_items), 6)
    item_area = card_h - 0.9
    spacing = min(item_area / max(n_before, 1), 0.55)
    y = card_top + 0.8
    for item in before_items[:n_before]:
        _add_accent_bar(slide, left_x + 0.2, y + 0.05, 0.04, 0.25, _t("accent5"))
        _add_text(slide, left_x + 0.4, y, left_w - 0.7, 0.4, item[:70],
                  font_size=FONT_SIZE_BODY, color=_t("body"))
        y += spacing

    # Vertical divider
    _add_vertical_line(slide, divider_x, card_top + 0.3, card_h - 0.6, 0.03, _t("border"))

    # Arrow in center
    arrow_y = card_top + card_h / 2 - 0.2
    arrow = slide.shapes.add_shape(
        MSO_SHAPE.RIGHT_ARROW, Inches(divider_x - 0.3), Inches(arrow_y),
        Inches(0.8), Inches(0.4)
    )
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = hex_to_rgb(_t("primary"))
    arrow.line.fill.background()

    # Right card: After (green tinted)
    _add_card_with_shadow(slide, right_x, card_top, right_w, card_h, _t("card_after"))
    _add_circle_icon(slide, right_x + 0.15, card_top + 0.15, 0.45, _t("accent4"), "O")
    _add_text(slide, right_x + 0.7, card_top + 0.18, right_w - 1.0, 0.4, "TO-BE (미래)",
              font_size=18, bold=True, color=_t("accent4"))

    n_after = min(len(after_items), 6)
    y = card_top + 0.8
    for item in after_items[:n_after]:
        _add_accent_bar(slide, right_x + 0.2, y + 0.05, 0.04, 0.25, _t("accent4"))
        _add_text(slide, right_x + 0.4, y, right_w - 0.7, 0.4, item[:70],
                  font_size=FONT_SIZE_BODY, color=_t("body"))
        y += spacing


def slide_kpi(prs, title, kpis):
    """KPI 카드 슬라이드 (Dashboard grid - mixed sizes + bar chart)."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, _t("background"))
    _slide_header(slide, "KEY PERFORMANCE INDICATORS", title)

    n = min(len(kpis), 4)
    if n == 0:
        return

    # Dashboard layout: left side cards, right side chart
    chart_area_w = 5.0
    card_area_w = CONTENT_W - chart_area_w - 0.4

    # Cards area (left) - mixed sizes: 1 large on top, rest smaller below
    card_x = MARGIN_X

    if n >= 1:
        # Large KPI card
        kpi0 = kpis[0]
        large_h = 2.2
        _add_card_with_shadow(slide, card_x, CONTENT_TOP, card_area_w, large_h, _t("card"))
        metric_name = kpi0.get("metric", kpi0.get("name", "KPI 1"))
        current = kpi0.get("current", "-")
        target = kpi0.get("target", "-")
        improvement = kpi0.get("improvement", "")

        _add_text(slide, card_x + 0.25, CONTENT_TOP + 0.15, card_area_w - 0.5, 0.3,
                  metric_name[:40], font_size=FONT_SIZE_SMALL, bold=True, color=_t("subtle"))
        _add_text(slide, card_x + 0.25, CONTENT_TOP + 0.55, card_area_w - 0.5, 0.7,
                  f"{current}  \u2192  {target}",
                  font_size=28, bold=True, color=_t("primary"),
                  font_name=_current_theme.get("font_title", "맑은 고딕"))
        if improvement:
            _add_rounded_card(slide, card_x + 0.25, CONTENT_TOP + large_h - 0.55, 2.2, 0.38,
                              _t("accent4"))
            _add_text(slide, card_x + 0.3, CONTENT_TOP + large_h - 0.53, 2.1, 0.34,
                      improvement[:25], font_size=FONT_SIZE_SMALL, bold=True,
                      color=_t("white"), alignment=PP_ALIGN.CENTER)

    # Smaller KPI cards below
    if n > 1:
        small_cols = min(n - 1, 3)
        small_cw = (card_area_w - (small_cols - 1) * 0.2) / small_cols
        small_ch = min(2.0, SAFE_BOTTOM - CONTENT_TOP - 2.5 - 0.2)
        small_y = CONTENT_TOP + 2.4

        for i, kpi in enumerate(kpis[1:n]):
            if i >= small_cols:
                break
            sx = card_x + i * (small_cw + 0.2)
            _add_card_with_shadow(slide, sx, small_y, small_cw, small_ch, _t("card"))

            metric_name = kpi.get("metric", kpi.get("name", f"KPI {i + 2}"))
            current = kpi.get("current", "-")
            target = kpi.get("target", "-")
            improvement = kpi.get("improvement", "")

            _add_text(slide, sx + 0.15, small_y + 0.1, small_cw - 0.3, 0.25,
                      metric_name[:25], font_size=10, bold=True, color=_t("subtle"))
            _add_text(slide, sx + 0.15, small_y + 0.4, small_cw - 0.3, 0.5,
                      f"{current} \u2192 {target}",
                      font_size=16, bold=True, color=_t("primary"),
                      font_name=_current_theme.get("font_title", "맑은 고딕"))
            if improvement:
                _add_rounded_card(slide, sx + 0.1, small_y + small_ch - 0.45, small_cw - 0.2,
                                  0.3, _t("accent4"))
                _add_text(slide, sx + 0.12, small_y + small_ch - 0.43, small_cw - 0.24, 0.26,
                          improvement[:20], font_size=9, bold=True,
                          color=_t("white"), alignment=PP_ALIGN.CENTER)

    # Bar chart on right side (current vs target comparison)
    chart_x = MARGIN_X + card_area_w + 0.4
    try:
        categories = []
        current_vals = []
        target_vals = []
        for kpi in kpis[:n]:
            name = kpi.get("metric", kpi.get("name", ""))[:15]
            categories.append(name)
            # Extract numeric values
            curr_str = str(kpi.get("current", "0"))
            tgt_str = str(kpi.get("target", "0"))
            curr_nums = re.findall(r'[\d.]+', curr_str)
            tgt_nums = re.findall(r'[\d.]+', tgt_str)
            current_vals.append(float(curr_nums[0]) if curr_nums else 0)
            target_vals.append(float(tgt_nums[0]) if tgt_nums else 0)

        if any(v > 0 for v in current_vals + target_vals):
            chart_buf = _create_bar_chart(
                categories, target_vals, _palette(), horizontal=True
            )
            if chart_buf:
                _insert_chart(slide, chart_buf, chart_x, CONTENT_TOP, chart_area_w, CONTENT_H)
    except Exception:
        pass  # Graceful degradation - skip chart on error


def slide_solution_highlight(prs, value_prop, overview):
    """솔루션 개요 하이라이트 (Full-width highlight)."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_gradient_bg(slide, _t("bg_section"), _t("gradient_end"), 'horizontal')

    _add_text(slide, 1.0, 0.6, 11, 0.4, "OUR SOLUTION",
              font_size=FONT_SIZE_METRIC_LABEL, bold=True, color=_t("muted"))

    # Large centered value proposition
    _add_text(slide, 1.0, 1.5, 11.3, 2.0, value_prop[:150],
              font_size=30, bold=True, color=_t("white"),
              font_name=_current_theme.get("font_title", "맑은 고딕"))

    if overview:
        _add_accent_bar(slide, 1.0, 3.8, 2.5, 0.05, _t("primary"))
        _add_text(slide, 1.0, 4.2, 11.3, 1.5, overview[:180],
                  font_size=17, color=_t("muted"))


def slide_scope(prs, in_scope, out_scope):
    """작업 범위 슬라이드 (Split 60/40 layout)."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, _t("background"))
    _slide_header(slide, "PROJECT SCOPE", "작업 범위")

    card_top = CONTENT_TOP
    card_h = SAFE_BOTTOM - card_top

    # 60/40 split
    left_w = CONTENT_W * 0.55
    right_w = CONTENT_W * 0.40
    left_x = MARGIN_X
    right_x = MARGIN_X + left_w + CONTENT_W * 0.05

    # In-Scope (larger, left)
    _add_card_with_shadow(slide, left_x, card_top, left_w, card_h, _t("card_after"))
    _add_circle_icon(slide, left_x + 0.15, card_top + 0.12, 0.38, _t("accent4"), "O")
    _add_text(slide, left_x + 0.65, card_top + 0.12, left_w - 1.0, 0.35,
              "포함 범위 (In-Scope)",
              font_size=15, bold=True, color=_t("accent4"))

    n_in = min(len(in_scope), 7)
    spacing = min(0.45, (card_h - 0.7) / max(n_in, 1))
    y = card_top + 0.65
    for item in in_scope[:n_in]:
        text = item if isinstance(item, str) else item.get("value", item.get("name", str(item)))
        _add_accent_bar(slide, left_x + 0.15, y + 0.05, 0.04, 0.25, _t("accent4"))
        _add_text(slide, left_x + 0.35, y, left_w - 0.6, 0.35, text[:60],
                  font_size=FONT_SIZE_BODY, color=_t("body"))
        y += spacing

    # Vertical divider
    divider_x = left_x + left_w + CONTENT_W * 0.025
    _add_vertical_line(slide, divider_x, card_top + 0.3, card_h - 0.6)

    # Out-of-Scope (smaller, right)
    _add_card_with_shadow(slide, right_x, card_top, right_w, card_h, _t("card_before"))
    _add_circle_icon(slide, right_x + 0.15, card_top + 0.12, 0.38, _t("accent5"), "X")
    _add_text(slide, right_x + 0.65, card_top + 0.12, right_w - 1.0, 0.35,
              "제외 범위 (Out of Scope)",
              font_size=15, bold=True, color=_t("accent5"))

    n_out = min(len(out_scope), 7)
    spacing_out = min(0.45, (card_h - 0.7) / max(n_out, 1))
    y = card_top + 0.65
    for item in out_scope[:n_out]:
        text = item if isinstance(item, str) else item.get("item", item.get("name", str(item)))
        _add_accent_bar(slide, right_x + 0.15, y + 0.05, 0.04, 0.25, _t("accent5"))
        _add_text(slide, right_x + 0.35, y, right_w - 0.6, 0.35, text[:55],
                  font_size=FONT_SIZE_BODY, color=_t("body"))
        y += spacing_out


def slide_features(prs, features):
    """핵심 기능 슬라이드 (카드 그리드, 와이드스크린)."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, _t("background"))
    _slide_header(slide, "KEY FEATURES", "핵심 기능", _t("secondary"), _t("secondary"))

    n = min(len(features), 6)
    cols = 3 if n > 4 else 2
    rows = math.ceil(n / cols)
    gap_x = 0.35
    gap_y = 0.3
    cw = (CONTENT_W - (cols - 1) * gap_x) / cols
    ch = min(2.2, (SAFE_BOTTOM - CONTENT_TOP - (rows - 1) * gap_y) / max(rows, 1))
    start_x = MARGIN_X
    palette = _palette()

    for i, feat in enumerate(features[:n]):
        col = i % cols
        row = i // cols
        x = start_x + col * (cw + gap_x)
        y = CONTENT_TOP + row * (ch + gap_y)

        _add_card_with_shadow(slide, x, y, cw, ch, _t("card_alt"))
        _add_circle_icon(slide, x + 0.15, y + 0.15, 0.38, palette[i % len(palette)],
                         str(i + 1))

        name = feat.get("name", f"기능 {i + 1}")
        desc = feat.get("description", "")
        _add_text(slide, x + 0.65, y + 0.15, cw - 0.85, 0.3, name[:30],
                  font_size=14, bold=True, color=_t("title"))
        _add_text(slide, x + 0.15, y + 0.6, cw - 0.3, ch - 0.75, desc[:90],
                  font_size=FONT_SIZE_SMALL, color=_t("body"))


def slide_tech_stack(prs, tech_items):
    """기술 스택 슬라이드."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, _t("background"))
    _slide_header(slide, "TECHNOLOGY STACK", "기술 스택", _t("accent3"), _t("accent3"))

    n = min(len(tech_items), 6)
    start_y, item_h = _calc_item_spacing(n, CONTENT_TOP, SAFE_BOTTOM, 0.78)
    palette = _palette()

    for i, item in enumerate(tech_items[:n]):
        y = start_y + i * item_h
        cat = item.get("category", "")
        tech = item.get("technology", item.get("tech", ""))

        _add_card_with_shadow(slide, MARGIN_X, y, CONTENT_W, item_h - 0.06, _t("card"))
        _add_circle_icon(slide, MARGIN_X + 0.15, y + (item_h - 0.42) / 2 - 0.03, 0.42,
                         palette[i % len(palette)], cat[:1].upper() if cat else "T")
        _add_text(slide, MARGIN_X + 0.75, y + 0.06, 3.0, 0.28, cat[:25],
                  font_size=FONT_SIZE_BODY, bold=True, color=_t("primary"))
        _add_text(slide, MARGIN_X + 0.75, y + 0.35, CONTENT_W - 1.2, 0.3, tech[:100],
                  font_size=14, color=_t("body"))


def slide_architecture(prs, diagram_path):
    """시스템 아키텍처 다이어그램 슬라이드."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, _t("background"))
    _slide_header(slide, "SYSTEM ARCHITECTURE", "시스템 아키텍처", _t("accent3"), _t("accent3"))

    img_top = CONTENT_TOP
    img_h = SAFE_BOTTOM - img_top
    img_w = CONTENT_W

    # Maintain image aspect ratio (original 1920x1080 = 16:9)
    aspect = 1920 / 1080
    fit_w = img_h * aspect
    if fit_w > img_w:
        fit_w = img_w
        fit_h = img_w / aspect
    else:
        fit_h = img_h

    img_left = MARGIN_X + (img_w - fit_w) / 2
    img_top_adj = img_top + (img_h - fit_h) / 2

    slide.shapes.add_picture(
        str(diagram_path),
        Inches(img_left), Inches(img_top_adj),
        Inches(fit_w), Inches(fit_h)
    )


def slide_timeline(prs, total_duration, phases):
    """타임라인 슬라이드 (Timeline/flow layout + Gantt chart)."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, _t("background"))
    _slide_header(slide, "PROJECT TIMELINE", "프로젝트 일정")

    # Duration badge
    if total_duration:
        badge_x = SLIDE_W - 3.0
        _add_rounded_card(slide, badge_x, 0.8, 2.3, 0.45, _t("primary"))
        _add_text(slide, badge_x + 0.05, 0.83, 2.2, 0.4, total_duration,
                  font_size=14, bold=True, color=_t("white"),
                  alignment=PP_ALIGN.CENTER)

    n = min(len(phases), 5)
    if n == 0:
        return

    palette = _palette()

    # Top half: Horizontal timeline/flow layout
    flow_y = CONTENT_TOP
    flow_h = 2.0
    node_size = 0.5
    node_gap = (CONTENT_W - n * node_size) / max(n, 1)
    node_gap = min(node_gap, 2.5)

    total_flow_w = n * node_size + (n - 1) * (node_gap - node_size) * 0.5
    start_x = MARGIN_X + (CONTENT_W - min(total_flow_w, CONTENT_W)) / 2
    actual_gap = (CONTENT_W - node_size) / max(n - 1, 1) if n > 1 else 0
    actual_gap = min(actual_gap, 2.8)

    for i, phase in enumerate(phases[:n]):
        x = start_x + i * actual_gap
        color = palette[i % len(palette)]
        name = phase.get("phase", phase.get("name", phase.get("phase_name", f"Phase {i + 1}")))
        duration = phase.get("duration", phase.get("period", ""))

        # Connecting line (before node, except first)
        if i > 0:
            prev_x = start_x + (i - 1) * actual_gap + node_size
            line_w = x - prev_x
            if line_w > 0:
                _add_accent_bar(slide, prev_x, flow_y + node_size / 2 - 0.02,
                                line_w, 0.04, _t("border"))

        # Node circle
        _add_circle_icon(slide, x, flow_y, node_size, color, str(i + 1))

        # Phase name below node
        _add_text(slide, x - 0.5, flow_y + node_size + 0.1, node_size + 1.0, 0.35,
                  name[:20], font_size=11, bold=True, color=_t("title"),
                  alignment=PP_ALIGN.CENTER)

        # Duration below name
        if duration:
            _add_text(slide, x - 0.5, flow_y + node_size + 0.4, node_size + 1.0, 0.3,
                      duration[:18], font_size=10, color=_t("subtle"),
                      alignment=PP_ALIGN.CENTER)

    # Bottom half: Gantt chart
    chart_y = CONTENT_TOP + flow_h + 1.2
    chart_h = SAFE_BOTTOM - chart_y
    if chart_h > 0.5:
        try:
            phase_names = [p.get("phase", p.get("name", f"P{i + 1}"))
                           for i, p in enumerate(phases[:n])]
            durations_list = [p.get("duration", p.get("period", "1"))
                              for p in phases[:n]]
            chart_buf = _create_timeline_chart(phase_names, durations_list)
            if chart_buf:
                _insert_chart(slide, chart_buf, MARGIN_X, chart_y, CONTENT_W, chart_h)
        except Exception:
            pass


def slide_team(prs, team_comp, total_mm):
    """팀 구성 슬라이드 (cards + donut chart)."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, _t("background"))
    _slide_header(slide, "PROJECT TEAM", "투입 인력", _t("secondary"), _t("secondary"))

    n = min(len(team_comp), 6)
    palette = _palette()

    # Layout: cards on left (60%), donut chart on right (40%)
    card_area_w = CONTENT_W * 0.58
    chart_area_w = CONTENT_W * 0.38
    chart_x = MARGIN_X + card_area_w + CONTENT_W * 0.04

    # Team cards (left side) - 2 columns
    cols = 2
    rows = math.ceil(n / cols)
    gap_x = 0.25
    gap_y = 0.25
    cw = (card_area_w - (cols - 1) * gap_x) / cols

    badge_h = 0.5
    badge_gap = 0.15
    available_h = SAFE_BOTTOM - CONTENT_TOP - badge_h - badge_gap
    ch = min(1.8, (available_h - (rows - 1) * gap_y) / max(rows, 1))
    start_x = MARGIN_X

    role_icons = ["PM", "BE", "FE", "UX", "QA", "DV"]

    for i, member in enumerate(team_comp[:n]):
        col = i % cols
        row = i // cols
        x = start_x + col * (cw + gap_x)
        y = CONTENT_TOP + row * (ch + gap_y)

        _add_card_with_shadow(slide, x, y, cw, ch, _t("card"))

        role = member.get("role", f"역할 {i + 1}")
        count = member.get("count", 1)
        expertise = member.get("expertise", member.get("skills", ""))

        icon_text = role_icons[i] if i < len(role_icons) else role[:2]
        _add_circle_icon(slide, x + (cw - 0.45) / 2, y + 0.1, 0.45,
                         palette[i % len(palette)], icon_text)

        _add_text(slide, x + 0.05, y + 0.62, cw - 0.1, 0.28, role[:18],
                  font_size=13, bold=True, color=_t("title"),
                  alignment=PP_ALIGN.CENTER)

        count_str = f"{count}명" if count >= 1 else f"{count} (파트타임)"
        _add_text(slide, x + 0.05, y + 0.9, cw - 0.1, 0.28, count_str,
                  font_size=18, bold=True, color=_t("primary"),
                  alignment=PP_ALIGN.CENTER,
                  font_name=_current_theme.get("font_title", "맑은 고딕"))

        if expertise and ch >= 1.5:
            _add_text(slide, x + 0.05, y + 1.2, cw - 0.1, 0.4,
                      str(expertise)[:40], font_size=9, color=_t("subtle"),
                      alignment=PP_ALIGN.CENTER)

    # Total effort badge
    badge_y = CONTENT_TOP + rows * (ch + gap_y) + 0.05
    badge_y = min(badge_y, SAFE_BOTTOM - badge_h)
    _add_rounded_card(slide, start_x, badge_y, card_area_w, badge_h, _t("primary"))
    _add_text(slide, start_x + 0.1, badge_y + 0.05, card_area_w - 0.2, 0.4,
              f"총 공수: {total_mm} Man-Months",
              font_size=16, bold=True, color=_t("white"),
              alignment=PP_ALIGN.CENTER,
              font_name=_current_theme.get("font_title", "맑은 고딕"))

    # Donut chart (right side) - team role distribution
    try:
        chart_data = []
        chart_colors = []
        for i, member in enumerate(team_comp[:n]):
            role = member.get("role", f"역할 {i + 1}")
            count = member.get("count", 1)
            chart_data.append((role[:15], max(count, 0.5)))
            chart_colors.append(palette[i % len(palette)])

        if chart_data:
            chart_buf = _create_donut_chart(chart_data, chart_colors)
            if chart_buf:
                chart_size = min(chart_area_w, CONTENT_H - 0.5)
                chart_cy = CONTENT_TOP + (CONTENT_H - chart_size) / 2
                _insert_chart(slide, chart_buf, chart_x, chart_cy, chart_size, chart_size)
    except Exception:
        pass


def slide_risks(prs, risks):
    """리스크 관리 슬라이드."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, _t("background"))
    _slide_header(slide, "RISK MANAGEMENT", "리스크 관리", _t("orange"), _t("orange"))

    n = min(len(risks), 4)
    start_y, item_h = _calc_item_spacing(n, CONTENT_TOP, SAFE_BOTTOM, 1.15)

    for i, risk in enumerate(risks[:n]):
        y = start_y + i * item_h
        risk_text = risk.get("risk", risk.get("description", ""))
        impact = risk.get("impact", risk.get("level", "MEDIUM"))
        mitigation = risk.get("mitigation", "")

        if impact == "HIGH":
            level_color = _t("accent5")
        elif impact == "LOW":
            level_color = _t("accent4")
        else:
            level_color = _t("amber")

        _add_card_with_shadow(slide, MARGIN_X, y, CONTENT_W, item_h - 0.08, _t("card"))

        # Impact badge
        _add_rounded_card(slide, MARGIN_X + 0.15, y + 0.12, 0.85, 0.32, level_color)
        _add_text(slide, MARGIN_X + 0.17, y + 0.13, 0.81, 0.3, impact,
                  font_size=10, bold=True, color=_t("white"),
                  alignment=PP_ALIGN.CENTER)

        _add_text(slide, MARGIN_X + 1.2, y + 0.08, CONTENT_W - 1.6, 0.35, risk_text[:70],
                  font_size=14, bold=True, color=_t("title"))

        if mitigation and item_h >= 0.85:
            _add_text(slide, MARGIN_X + 1.2, y + 0.5, CONTENT_W - 1.6, 0.4,
                      f"\u2192 {mitigation[:80]}",
                      font_size=FONT_SIZE_SMALL, color=_t("subtle"))


def slide_benefits(prs, quant_benefits, qual_benefits):
    """기대 효과 슬라이드 (Dashboard grid + bar chart)."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, _t("background"))
    _slide_header(slide, "EXPECTED BENEFITS", "기대 효과", _t("accent4"), _t("accent4"))

    has_quant = bool(quant_benefits)
    has_qual = bool(qual_benefits)

    # Dashboard layout: large card left, chart right
    if has_quant:
        n_q = min(len(quant_benefits), 4)

        # Mixed-size dashboard: 1 large + smaller cards
        if n_q == 1:
            # Single large card
            b = quant_benefits[0]
            _add_card_with_shadow(slide, MARGIN_X, CONTENT_TOP, CONTENT_W, 2.0, _t("card"))
            metric = b.get("metric", "")
            before = b.get("before", "-")
            after = b.get("after", "-")
            _add_text(slide, MARGIN_X + 0.3, CONTENT_TOP + 0.15, CONTENT_W - 0.6, 0.3,
                      metric[:40], font_size=13, bold=True, color=_t("subtle"))
            _add_text(slide, MARGIN_X + 0.3, CONTENT_TOP + 0.55, CONTENT_W - 0.6, 0.7,
                      f"{before}  \u2192  {after}",
                      font_size=32, bold=True, color=_t("accent4"),
                      alignment=PP_ALIGN.CENTER,
                      font_name=_current_theme.get("font_title", "맑은 고딕"))
        else:
            # 2 medium cards on left, chart on right
            card_w = (CONTENT_W * 0.55 - 0.3) / 2
            chart_w = CONTENT_W * 0.40

            for i, b in enumerate(quant_benefits[:min(n_q, 4)]):
                col = i % 2
                row = i // 2
                x = MARGIN_X + col * (card_w + 0.3)
                y = CONTENT_TOP + row * 1.8
                _add_card_with_shadow(slide, x, y, card_w, 1.6, _t("card"))

                metric = b.get("metric", "")
                before = b.get("before", "-")
                after = b.get("after", "-")

                _add_text(slide, x + 0.15, y + 0.1, card_w - 0.3, 0.25,
                          metric[:30], font_size=FONT_SIZE_SMALL, bold=True,
                          color=_t("subtle"))
                _add_text(slide, x + 0.15, y + 0.45, card_w - 0.3, 0.55,
                          f"{before} \u2192 {after}",
                          font_size=18, bold=True, color=_t("accent4"),
                          alignment=PP_ALIGN.CENTER,
                          font_name=_current_theme.get("font_title", "맑은 고딕"))
                _add_circle_icon(slide, x + (card_w - 0.3) / 2, y + 1.15,
                                 0.3, _t("accent4"), "\u2191")

            # Bar chart on right
            try:
                categories = [b.get("metric", "")[:15] for b in quant_benefits[:n_q]]
                values = []
                for b in quant_benefits[:n_q]:
                    after_str = str(b.get("after", "0"))
                    nums = re.findall(r'[\d.]+', after_str)
                    values.append(float(nums[0]) if nums else 0)

                if any(v > 0 for v in values):
                    chart_buf = _create_bar_chart(categories, values, _palette(),
                                                  horizontal=True)
                    if chart_buf:
                        chart_x = MARGIN_X + CONTENT_W * 0.58
                        _insert_chart(slide, chart_buf, chart_x, CONTENT_TOP,
                                      chart_w, min(CONTENT_H, 3.5))
            except Exception:
                pass

    qual_start_y = CONTENT_TOP + (2.0 if has_quant and len(quant_benefits) <= 2 else
                                   3.8 if has_quant else 0)
    if has_qual and qual_start_y < SAFE_BOTTOM - 0.5:
        n_qual = min(len(qual_benefits), 5)
        available = SAFE_BOTTOM - qual_start_y
        spacing = min(0.42, available / max(n_qual, 1))
        y = qual_start_y
        for i, qual in enumerate(qual_benefits[:n_qual]):
            if y + 0.35 > SAFE_BOTTOM:
                break
            text = qual if isinstance(qual, str) else str(qual)
            _add_accent_bar(slide, MARGIN_X + 0.1, y + 0.08, 0.04, 0.2, _t("accent4"))
            _add_text(slide, MARGIN_X + 0.3, y, CONTENT_W - 0.5, 0.35,
                      text[:90], font_size=14, color=_t("body"))
            y += spacing


def slide_next_steps(prs, steps):
    """다음 단계 슬라이드 (Timeline/flow layout)."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, _t("background"))
    _slide_header(slide, "NEXT STEPS", "다음 단계")

    n = min(len(steps), 5)
    if n == 0:
        return

    palette = _palette()

    # Horizontal flow layout at top
    flow_y = CONTENT_TOP
    node_size = 0.55
    actual_gap = min((CONTENT_W - node_size) / max(n - 1, 1), 2.8) if n > 1 else 0
    start_x = MARGIN_X + (CONTENT_W - (node_size + (n - 1) * actual_gap)) / 2

    for i, step in enumerate(steps[:n]):
        x = start_x + i * actual_gap
        color = palette[i % len(palette)]
        action = step.get("action", step) if isinstance(step, dict) else str(step)
        duration = step.get("duration", "") if isinstance(step, dict) else ""

        # Connecting line
        if i > 0:
            prev_x = start_x + (i - 1) * actual_gap + node_size
            line_w = x - prev_x
            if line_w > 0:
                _add_accent_bar(slide, prev_x, flow_y + node_size / 2 - 0.02,
                                line_w, 0.04, _t("border"))

        # Node
        _add_circle_icon(slide, x, flow_y, node_size, color, str(i + 1))

        # Card below node
        card_y = flow_y + node_size + 0.2
        card_w = actual_gap - 0.15 if actual_gap > 0 else 2.5
        card_w = min(card_w, 2.8)
        card_h = min(2.8, SAFE_BOTTOM - card_y - 0.1)

        card_x = x + node_size / 2 - card_w / 2
        card_x = max(MARGIN_X, min(card_x, SLIDE_W - MARGIN_X - card_w))

        _add_card_with_shadow(slide, card_x, card_y, card_w, card_h, _t("card"))
        _add_text(slide, card_x + 0.1, card_y + 0.1, card_w - 0.2, card_h - 0.5,
                  action[:80], font_size=FONT_SIZE_SMALL, color=_t("body"))
        if duration:
            _add_text(slide, card_x + 0.1, card_y + card_h - 0.4, card_w - 0.2, 0.3,
                      duration[:18], font_size=10, bold=True, color=_t("subtle"),
                      alignment=PP_ALIGN.CENTER)


def slide_closing(prs, title_text, cta):
    """마무리 슬라이드 (그라데이션)."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_gradient_bg(slide, _t("bg_section"), _t("gradient_end"), 'diagonal')

    _add_text(slide, 0.5, 2.0, SLIDE_W - 1, 1.2, title_text,
              font_size=48, bold=True, color=_t("white"),
              alignment=PP_ALIGN.CENTER,
              font_name=_current_theme.get("font_title", "맑은 고딕"))
    _add_accent_bar(slide, SLIDE_W / 2 - 1.5, 3.5, 3.0, 0.05, _t("primary"))
    _add_text(slide, 0.5, 3.9, SLIDE_W - 1, 0.6, cta,
              font_size=19, color=_t("muted"), alignment=PP_ALIGN.CENTER)
    _add_text(slide, 0.5, 5.0, SLIDE_W - 1, 0.5, "Q & A",
              font_size=28, bold=True, color=_t("primary"),
              alignment=PP_ALIGN.CENTER,
              font_name=_current_theme.get("font_title", "맑은 고딕"))


# ============================================================
# Validation
# ============================================================

def validate_ppt(prs):
    """PPT 슬라이드별 레이아웃 검증. 오버플로우/겹침 검출."""
    slide_w = prs.slide_width
    slide_h = prs.slide_height
    issues = []

    for idx, slide in enumerate(prs.slides):
        slide_num = idx + 1
        shapes_info = []

        for shape in slide.shapes:
            left = shape.left / 914400
            top = shape.top / 914400
            width = shape.width / 914400
            height = shape.height / 914400
            right = left + width
            bottom = top + height

            shapes_info.append({
                "name": shape.shape_id,
                "left": round(left, 2),
                "top": round(top, 2),
                "right": round(right, 2),
                "bottom": round(bottom, 2),
            })

            slide_w_in = slide_w / 914400
            slide_h_in = slide_h / 914400

            if right > slide_w_in + 0.1:
                issues.append(f"  슬라이드 {slide_num}: 요소가 우측 경계 초과 "
                              f"(right={right:.2f} > {slide_w_in:.1f})")
            if bottom > slide_h_in + 0.1:
                issues.append(f"  슬라이드 {slide_num}: 요소가 하단 경계 초과 "
                              f"(bottom={bottom:.2f} > {slide_h_in:.1f})")

        for i, a in enumerate(shapes_info):
            for j, b in enumerate(shapes_info):
                if i >= j:
                    continue
                same_pos = abs(a["left"] - b["left"]) < 0.05 and abs(a["top"] - b["top"]) < 0.05
                same_size = abs(a["right"] - a["left"] - (b["right"] - b["left"])) < 0.1 and \
                            abs(a["bottom"] - a["top"] - (b["bottom"] - b["top"])) < 0.1
                if same_pos and same_size:
                    a_area = (a["right"] - a["left"]) * (a["bottom"] - a["top"])
                    if a_area > 1.0:
                        issues.append(
                            f"  슬라이드 {slide_num}: 동일 위치 중복 요소 "
                            f"(top={a['top']:.1f}, left={a['left']:.1f})"
                        )

    return issues


# ============================================================
# Data Normalization
# ============================================================

def normalize_proposal_data(data: dict) -> dict:
    """제안서 JSON을 PPT 데이터로 정규화."""
    if "storytelling_structure" in data:
        return data
    if "current_challenges" in data:
        return _normalize_new_format(data)
    return _normalize_legacy_format(data)


def _normalize_new_format(data: dict) -> dict:
    """새 제안서 JSON 포맷 정규화."""
    exec_summary = data.get("executive_summary", {})
    core_msg = exec_summary.get("core_message", "")
    invest = exec_summary.get("investment_overview", {})
    key_metrics = exec_summary.get("key_metrics", {})
    title = data.get("title", "프로젝트 제안서")

    normalized = {
        "title": title,
        "metadata": {
            "proposal_date": data.get("created_at", ""),
            "client_company": data.get("client_name", ""),
            "proposer": data.get("proposer", ""),
        },
        "storytelling_structure": {
            "hook": core_msg[:60] if core_msg else title,
            "solution": core_msg[:80] if core_msg else title,
            "cta": "함께 시작하겠습니다",
        },
    }

    normalized["executive_summary"] = {
        "problem": core_msg,
        "solution": core_msg,
        "duration": invest.get("project_period", ""),
        "effort": invest.get("total_effort_with_buffer", invest.get("total_effort", "")),
        "key_benefits": list(key_metrics.values())[:2] if key_metrics else [],
    }

    challenges_raw = data.get("current_challenges", [])
    challenges = [{"area": ch.get("title", ""), "symptom": ch.get("description", ""),
                    "business_impact": ch.get("description", "")} for ch in challenges_raw]

    goals = data.get("project_goals", [])
    future_vision = {f"vision_{i+1}": f"{g.get('title', '')}: {g.get('description', '')}"
                     for i, g in enumerate(goals[:4])}

    normalized["current_situation"] = {
        "challenges": challenges,
        "risks_if_no_change": [ch.get("description", "") for ch in challenges_raw[:4]],
        "future_vision": future_vision,
    }

    kpis = [{"metric": k.get("name", ""), "current": k.get("current", "-"),
             "target": k.get("target", ""), "improvement": k.get("measurement", "")}
            for k in data.get("kpi", [])[:4]]
    normalized["objectives"] = {"kpis": kpis, "goals": [g.get("description", "") for g in goals]}

    scope_data = data.get("scope", {})
    in_scope = [{"value": s} if isinstance(s, str) else s for s in scope_data.get("in_scope", [])]
    out_scope = [{"item": s} if isinstance(s, str) else s for s in scope_data.get("out_of_scope", [])]
    modules = data.get("solution", {}).get("modules", [])

    normalized["solution"] = {
        "value_proposition": core_msg[:60] if core_msg else title,
        "overview": ", ".join(m.get("name", "") for m in modules[:3]),
        "scope": {"in_scope": in_scope, "out_of_scope": out_scope},
    }

    tech_data = data.get("technology_stack", {})
    tech_list = [{"category": k, "technology": ", ".join(v.values()) if isinstance(v, dict) else str(v)}
                 for k, v in tech_data.items() if isinstance(v, dict)]
    normalized["technical_approach"] = {"technology_stack": tech_list}

    timeline_data = data.get("timeline", {})
    normalized["timeline"] = {
        "total_duration": timeline_data.get("total_duration", ""),
        "phases": [{"phase": p.get("name", ""), "duration": p.get("period", ""),
                     "period": p.get("dates", "")} for p in timeline_data.get("phases", [])],
    }

    resource = data.get("resource_plan", {})
    normalized["team"] = {
        "composition": [{"role": m.get("role", ""), "count": m.get("count", 1),
                          "expertise": m.get("period", "")}
                         for m in resource.get("team", [])],
        "effort_summary": {"total": {"man_months": resource.get("total_man_months_with_buffer",
                                                                  resource.get("total_man_months", 0))}},
    }

    normalized["risk_management"] = [{"risk": r.get("title", r.get("description", "")),
                                       "impact": r.get("impact", "MEDIUM"),
                                       "mitigation": r.get("mitigation", "")}
                                      for r in data.get("risks", [])]

    benefits = data.get("expected_benefits", {})
    if isinstance(benefits, dict):
        quant = [{"metric": b.get("item", ""), "before": b.get("before", ""),
                  "after": b.get("after", ""), "improvement": "개선"}
                 for b in benefits.get("quantitative", [])[:4]]
        normalized["expected_benefits"] = {"quantitative": quant, "qualitative": benefits.get("qualitative", [])}
    else:
        normalized["expected_benefits"] = {"quantitative": [], "qualitative": []}

    milestones = timeline_data.get("milestones", [])
    normalized["next_steps"] = [{"step": i + 1, "action": ms.get("name", ""), "duration": ms.get("date", "")}
                                 for i, ms in enumerate(milestones[:5])]
    return normalized


def _normalize_legacy_format(data: dict) -> dict:
    """기존 ProposalDocument 모델 JSON 정규화."""
    title = data.get("title", "프로젝트 제안서")
    exec_str = data.get("executive_summary", "")
    overview = data.get("project_overview", {})
    timeline_data = data.get("timeline", {})
    resource = data.get("resource_plan", {})
    total_duration = timeline_data.get("total_duration", "")
    total_mm = resource.get("total_man_months", "")

    normalized = {
        "title": title,
        "metadata": data.get("metadata", {}),
        "storytelling_structure": {
            "hook": exec_str[:60] if isinstance(exec_str, str) and exec_str else title,
            "solution": exec_str[:80] if isinstance(exec_str, str) and exec_str else title,
            "cta": "함께 시작하겠습니다",
        },
    }

    if isinstance(exec_str, str):
        normalized["executive_summary"] = {
            "problem": exec_str,
            "solution": exec_str,
            "duration": total_duration,
            "effort": f"{total_mm} M/M" if total_mm else "",
            "key_benefits": data.get("expected_benefits", [])[:2]
                            if isinstance(data.get("expected_benefits"), list) else []
        }
    else:
        normalized["executive_summary"] = exec_str

    background = overview.get("background", "")
    objectives = overview.get("objectives", [])
    success_criteria = overview.get("success_criteria", [])

    challenges = []
    if background:
        sentences = [s.strip() for s in background.replace('. ', '.\n').split('\n') if s.strip()]
        for i, sent in enumerate(sentences[:4]):
            challenges.append({"area": f"현황 {i+1}", "symptom": sent[:80], "business_impact": sent[:80]})

    risks_if_no = [f"{obj} 미달성 시 서비스 품질 저하" for obj in objectives[:4] if isinstance(obj, str)]
    future_vision = {f"vision_{i+1}": c for i, c in enumerate(success_criteria[:4])}

    normalized["current_situation"] = {
        "challenges": challenges,
        "risks_if_no_change": risks_if_no if risks_if_no else ["현재 방식 유지 시 비효율 지속"],
        "future_vision": future_vision,
    }

    kpis = [{"metric": c[:30], "current": "-", "target": c, "improvement": ""}
            for c in success_criteria[:4]]
    normalized["objectives"] = {"kpis": kpis, "goals": objectives}

    scope = data.get("scope_of_work", {})
    solution_approach = data.get("solution_approach", {})
    in_scope_list = scope.get("in_scope", [])
    out_scope_list = scope.get("out_of_scope", [])
    key_features = scope.get("key_features", [])

    in_scope_converted = [{"value": i} if isinstance(i, str) else i for i in in_scope_list]
    out_scope_converted = [{"item": i} if isinstance(i, str) else i for i in out_scope_list]

    normalized["solution"] = {
        "value_proposition": solution_approach.get("overview", title),
        "overview": solution_approach.get("architecture", ""),
        "scope": {"in_scope": in_scope_converted, "out_of_scope": out_scope_converted},
        "key_features": key_features,
    }

    tech_stack = solution_approach.get("technology_stack", [])
    tech_converted = []
    for item in tech_stack:
        if isinstance(item, str):
            parts = item.split("(", 1)
            cat = parts[1].rstrip(")") if len(parts) > 1 else ""
            tech_converted.append({"category": cat, "technology": parts[0].strip()})
        else:
            tech_converted.append(item)
    normalized["technical_approach"] = {"technology_stack": tech_converted}

    phases = timeline_data.get("phases", [])
    normalized["timeline"] = {
        "total_duration": total_duration,
        "phases": [{"phase": p.get("phase_name", ""), "duration": p.get("duration", ""),
                     "period": p.get("description", "")} for p in phases],
    }

    team_structure = resource.get("team_structure", [])
    normalized["team"] = {
        "composition": [{"role": m.get("role", ""), "count": m.get("count", 1),
                          "expertise": ", ".join(m.get("responsibilities", [])[:2])}
                         for m in team_structure],
        "effort_summary": {"total": {"man_months": total_mm}},
    }

    normalized["risk_management"] = [{"risk": r.get("description", ""),
                                       "impact": r.get("level", "MEDIUM"),
                                       "mitigation": r.get("mitigation", "")}
                                      for r in data.get("risks", [])]

    benefits = data.get("expected_benefits", [])
    if isinstance(benefits, list):
        quant = [{"metric": b[:40], "before": "", "after": "", "improvement": ""}
                 for b in benefits[:4] if isinstance(b, str)]
        normalized["expected_benefits"] = {"quantitative": quant, "qualitative": benefits}
    else:
        normalized["expected_benefits"] = benefits

    steps = data.get("next_steps", [])
    normalized["next_steps"] = [{"step": i + 1, "action": s, "duration": ""}
                                 if isinstance(s, str) else s for i, s in enumerate(steps)]
    return normalized


# ============================================================
# Main Generator
# ============================================================

def _get_or_generate_arch_diagram() -> Path | None:
    """아키텍처 다이어그램 PNG 가져오기 또는 TRD에서 자동 생성."""
    project_root = Path(__file__).parent.parent.parent
    diagram_dir = project_root / "workspace/outputs/diagrams"
    existing = list(diagram_dir.glob("ARCH-*.png")) if diagram_dir.exists() else []
    if existing:
        return max(existing, key=lambda x: x.stat().st_mtime)

    trd_dir = project_root / "workspace/outputs/trd"
    trd_files = list(trd_dir.glob("TRD-*.json")) if trd_dir.exists() else []
    if not trd_files:
        return None

    trd_path = max(trd_files, key=lambda x: x.stat().st_mtime)
    try:
        print(f"   아키텍처 다이어그램 자동 생성 중... ({trd_path.name})")
        result = generate_from_trd_file(trd_path, diagram_dir)
        print(f"   다이어그램 생성 완료: {result.name}")
        return result
    except Exception as e:
        print(f"   [경고] 다이어그램 생성 실패: {e}")
        return None


def load_proposal_json(json_path: Path) -> dict:
    with open(json_path, 'r', encoding='utf-8') as f:
        return json.load(f)


def generate_ppt(proposal_path: Path, output_path: Path, theme_name: str = "modern-blue"):
    """PPT 생성 메인."""
    global _current_theme
    _current_theme = PPT_THEMES.get(theme_name, PPT_THEMES["modern-blue"])

    json_path = proposal_path.with_suffix('.json')
    if json_path.exists():
        raw_data = load_proposal_json(json_path)
        print(f"   JSON 데이터 로드: {json_path.name}")
    else:
        raw_data = {"title": "프로젝트 제안서"}

    data = normalize_proposal_data(raw_data)

    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)

    title = data.get("title", "프로젝트 제안서").replace(" 제안서", "")
    metadata = data.get("metadata", {})
    exec_summary = data.get("executive_summary", {})
    current_sit = data.get("current_situation", {})
    objectives = data.get("objectives", {})
    solution = data.get("solution", {})
    tech = data.get("technical_approach", {})
    timeline = data.get("timeline", {})
    team = data.get("team", {})
    risks = data.get("risk_management", [])
    benefits = data.get("expected_benefits", {})
    next_steps = data.get("next_steps", [])
    storytelling = data.get("storytelling_structure", {})
    key_features = solution.get("key_features", [])

    client = metadata.get("client_company", data.get("client_name", ""))
    date_str = metadata.get("proposal_date", metadata.get("created_at", ""))
    if not date_str:
        date_str = datetime.now().strftime("%Y-%m-%d")
    subtitle = f"{client} 귀중" if client else ""

    total_duration = timeline.get("total_duration", "")
    total_mm = team.get("effort_summary", {}).get("total", {}).get("man_months", "")

    # ---- 슬라이드 생성 ----

    slide_cover(prs, title, subtitle, date_str)

    toc_items = [
        "경영진 요약 (Executive Summary)",
        "현재 상황과 도전 (Challenges)",
        "프로젝트 목표 (Objectives)",
        "우리의 솔루션 (Solution)",
        "핵심 기능 (Features)",
        "기술 스택 (Tech Stack)",
        "시스템 아키텍처 (Architecture)",
        "프로젝트 일정 (Timeline)",
        "투입 인력 (Team)",
        "리스크 관리 (Risks)",
        "기대 효과 (Benefits)",
        "다음 단계 (Next Steps)",
    ]
    slide_toc(prs, toc_items)

    # 경영진 요약 - 핵심
    exec_problem = exec_summary.get("problem", "") if isinstance(exec_summary, dict) else str(exec_summary)
    exec_metrics = []
    if total_duration:
        exec_metrics.append({"label": "프로젝트 기간", "value": total_duration, "desc": ""})
    if total_mm:
        exec_metrics.append({"label": "투입 공수", "value": f"{total_mm} M/M", "desc": ""})
    key_benefits = exec_summary.get("key_benefits", []) if isinstance(exec_summary, dict) else []
    if key_benefits:
        exec_metrics.append({"label": "핵심 효과", "value": str(key_benefits[0])[:20], "desc": ""})
    slide_exec_highlight(prs, exec_problem[:150], exec_metrics)

    # 경영진 요약 - 상세
    detail_items = []
    if isinstance(exec_summary, dict):
        if exec_summary.get("problem"):
            detail_items.append({"label": "해결할 문제", "value": exec_summary["problem"][:80]})
        if exec_summary.get("solution") and exec_summary["solution"] != exec_summary.get("problem"):
            detail_items.append({"label": "솔루션", "value": exec_summary["solution"][:80]})
        if total_duration:
            detail_items.append({"label": "예상 기간", "value": total_duration})
        if exec_summary.get("effort"):
            detail_items.append({"label": "투입 공수", "value": str(exec_summary["effort"])})
        for b in key_benefits[:2]:
            detail_items.append({"label": "핵심 효과", "value": str(b)[:60]})
    slide_exec_detail(prs, detail_items)

    # 섹션: 현재 상황
    slide_section(prs, 1, "현재 상황", "Current Challenges")

    challenges = current_sit.get("challenges", [])
    slide_challenges(prs, "현재의 도전과 과제", challenges)

    risks_no_change = current_sit.get("risks_if_no_change", [])
    slide_risks_no_change(prs, risks_no_change)

    before_items = [c.get("business_impact", c.get("symptom", "")) for c in challenges[:5]]
    after_items = list(current_sit.get("future_vision", {}).values())[:5]
    slide_before_after(prs, before_items, after_items)

    # 섹션: 프로젝트 목표
    slide_section(prs, 2, "프로젝트 목표", "Objectives & KPI")
    kpis = objectives.get("kpis", [])
    slide_kpi(prs, "핵심 성과 지표 (KPI)", kpis)

    # 섹션: 솔루션
    slide_section(prs, 3, "우리의 솔루션", "Our Solution")
    slide_solution_highlight(prs, solution.get("value_proposition", title), solution.get("overview", "")[:120])

    in_scope = solution.get("scope", {}).get("in_scope", [])
    out_scope = solution.get("scope", {}).get("out_of_scope", [])
    slide_scope(prs, in_scope, out_scope)

    if key_features:
        slide_features(prs, key_features)

    # 기술 스택
    tech_stack = tech.get("technology_stack", [])
    slide_tech_stack(prs, tech_stack)

    # 시스템 아키텍처 다이어그램
    arch_diagram_path = _get_or_generate_arch_diagram()
    if arch_diagram_path:
        slide_architecture(prs, arch_diagram_path)

    # 섹션: 일정
    slide_section(prs, 4, "프로젝트 일정", "Timeline & Milestones")
    phases = timeline.get("phases", [])
    slide_timeline(prs, total_duration, phases)

    # 팀 구성
    team_comp = team.get("composition", [])
    slide_team(prs, team_comp, total_mm)

    # 리스크
    slide_risks(prs, risks)

    # 기대 효과
    quant = benefits.get("quantitative", []) if isinstance(benefits, dict) else []
    qual = benefits.get("qualitative", []) if isinstance(benefits, dict) else (
        benefits if isinstance(benefits, list) else [])
    slide_benefits(prs, quant, qual)

    # 다음 단계
    slide_next_steps(prs, next_steps)

    # Q&A
    slide_closing(prs, "감사합니다", storytelling.get("cta", "함께 시작하겠습니다"))

    # ---- 최종 검증 ----
    print("\n[검증] 슬라이드 레이아웃 검증 중...")
    issues = validate_ppt(prs)
    if issues:
        print(f"  [주의] {len(issues)}건의 레이아웃 이슈 발견:")
        for issue in issues[:10]:
            print(f"    {issue}")
    else:
        print("  [통과] 모든 슬라이드 레이아웃 정상")

    prs.save(str(output_path))
    slide_count = len(prs.slides)
    return output_path, slide_count


def main():
    parser = argparse.ArgumentParser(description="PPT 제안서 생성")
    parser.add_argument("--theme", choices=list(PPT_THEMES.keys()),
                        default="modern-blue",
                        help="PPT 테마 선택 (default: modern-blue)")
    args = parser.parse_args()

    print("\n" + "=" * 70)
    print("PPT 제안서 생성 (2026 Premium Design)")
    print(f'시작 시간: {datetime.now().strftime("%Y-%m-%d %H:%M:%S")}')
    print(f'테마: {args.theme}')
    print("=" * 70)

    project_root = Path(__file__).parent.parent.parent
    proposal_dir = project_root / "workspace/outputs/proposals"
    md_files = list(proposal_dir.glob("PROP-*.md"))

    if not md_files:
        print("제안서 파일을 찾을 수 없습니다.")
        print("먼저 /pro:pro-maker를 실행하세요.")
        return

    proposal_path = max(md_files, key=lambda x: x.stat().st_mtime)
    print(f"\n[입력] 제안서: {proposal_path}")

    output_dir = project_root / "workspace/outputs/ppt"
    output_dir.mkdir(parents=True, exist_ok=True)

    timestamp = datetime.now().strftime("%Y%m%d-%H%M%S")
    output_path = output_dir / f"PPT-{timestamp}.pptx"

    result, slide_count = generate_ppt(proposal_path, output_path, theme_name=args.theme)

    if result:
        print(f"\n[완료] PPT 생성 완료: {output_path}")
        print(f"   슬라이드 수: {slide_count}장")
        print(f"   테마: {args.theme}")
        print(f"   해상도: 16:9 와이드스크린 ({SLIDE_W}\" x {SLIDE_H}\")")
    else:
        print("\n[실패] PPT 생성 실패")


if __name__ == "__main__":
    main()
